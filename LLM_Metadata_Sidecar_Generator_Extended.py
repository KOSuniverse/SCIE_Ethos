"""
LLM Metadata Sidecar Generator Script (Extended)
-------------------------------------------------
This script scans a project directory recursively for supported document types,
extracts structured metadata from the document itself (if present),
or falls back to OpenAI GPT to generate metadata.
It then saves a sidecar .json file into a `_metadata/` subfolder next to each file.

Supported: .docx, .xlsx, .pptx, .pdf
"""

import json
from datetime import datetime
from docx import Document
import openpyxl
from pptx import Presentation
import pdfplumber
from openai import OpenAI
import numpy as np
import pandas as pd
import streamlit as st
import os

# --- Local-only global alias helpers ---
def load_global_aliases():
    alias_path = os.path.join(os.getcwd(), "global_column_aliases.json")
    if os.path.exists(alias_path):
        with open(alias_path, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}

def update_global_aliases(aliases):
    alias_path = os.path.join(os.getcwd(), "global_column_aliases.json")
    with open(alias_path, "w", encoding="utf-8") as f:
        json.dump(aliases, f, indent=2)

# --- Local file scanning and metadata writing ---
SUPPORTED_EXTS = (".xlsx", ".docx", ".pdf", ".pptx")

def find_all_supported_files(base_dir):
    return [
        os.path.join(root, f)
        for root, dirs, files in os.walk(base_dir)
        for f in files if f.lower().endswith(SUPPORTED_EXTS)
    ]

def extract_text_for_metadata(path):
    ext = path.lower().split('.')[-1]
    try:
        if ext == "docx":
            doc = Document(path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return text
        elif ext == "xlsx":
            wb = openpyxl.load_workbook(path, read_only=True)
            text = []
            for sheet in wb.worksheets:
                text.append(f"Sheet: {sheet.title}")
                headers = [cell.value for cell in next(sheet.iter_rows(max_row=1))]
                if headers:
                    text.append("Columns: " + " | ".join([str(h) for h in headers if h]))
            return "\n".join(text)
        elif ext == "pptx":
            prs = Presentation(path)
            text = "\n".join(
                shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame
            )
            return text
        elif ext == "pdf":
            with pdfplumber.open(path) as pdf:
                extracted_text = [page.extract_text() for page in pdf.pages if page.extract_text()]
                if any(extracted_text):
                    return "\n".join(extracted_text)
    except Exception as e:
        st.warning(f"Failed to extract text from {path}: {e}")
    return ""

def auto_generate_metadata(text):
    prompt = f"""
You are an expert document classification system trained to extract structured metadata from internal business documents.\n\nYour task is to analyze the following file content and return metadata that will be used to:\n- Select relevant documents in response to business questions\n- Embed file context for AI search\n- Categorize content for analytics and modeling\n\n📝 Analyze the content below (first 3000 characters):\n{text[:3000]}\n\n🎯 Return structured JSON metadata in the format:\n{{\n  \"title\": \"Short descriptive title (10 words max)\",\n  \"summary\": \"Detailed overview of the document content, including its business purpose, key fields, sheet/section names, and any key metrics or concepts. Use 4–6 sentences or more if needed.\",\n  \"tags\": [\n    \"specific keyword or metric\",\n    \"synonym or variant if applicable\",\n    \"limit 7–10 tags\"\n  ],\n  \"category\": \"Primary topic or department (e.g., 'supply_chain', 'finance', 'compliance')\"\n}}\n\nOutput only valid JSON."
    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2
        )
        raw = response.choices[0].message.content.strip()
        if raw.startswith("```json") or raw.startswith("```"):
            raw = raw.strip("`").replace("json", "").strip()
        return json.loads(raw)
    except Exception as e:
        st.warning(f"LLM metadata generation failed: {e}")
        return {}

def scan_and_process_local(base_dir):
    st.info(f"Scanning {base_dir} for supported files...")
    global_aliases = load_global_aliases()
    updated_global_aliases = global_aliases.copy()
    files = find_all_supported_files(base_dir)
    saved_files = []
    for file_path in files:
        ext = file_path.split(".")[-1].lower()
        meta = {}
        text = extract_text_for_metadata(file_path)
        gpt_meta = auto_generate_metadata(text)
        for key in ["title", "category", "tags", "summary"]:
            meta[key] = gpt_meta.get(key, "")
        meta["source_file"] = file_path
        meta["last_indexed"] = datetime.now().isoformat()
        meta["last_modified"] = datetime.fromtimestamp(os.path.getmtime(file_path)).isoformat() if os.path.exists(file_path) else None
        if ext == "xlsx":
            try:
                profile_data = profile_excel(file_path)
            except Exception as e:
                st.warning(f"Profile data error: {e}")
                profile_data = {}
            meta["profile_data"] = profile_data
            all_columns = []
            for sheet in profile_data.values():
                all_columns.extend(list(sheet.keys()))
            meta["columns"] = list(set(all_columns))
            column_aliases = map_columns_to_concepts(meta["columns"], updated_global_aliases)
            meta["column_aliases"] = column_aliases
            updated_global_aliases.update(column_aliases)
        summary_text = meta.get("summary", "")
        embedding = get_embedding(summary_text) if summary_text else None
        meta["embedding"] = embedding.tolist() if embedding is not None else []
        meta_path = write_sidecar_metadata(file_path, meta)
        st.success(f"✅ Metadata saved: {meta_path}")
        saved_files.append(meta_path)
    update_global_aliases(updated_global_aliases)
    st.info(f"Finished. {len(saved_files)} metadata files written.")


client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

def get_embedding(text):
    response = client.embeddings.create(
        model="text-embedding-ada-002",
        input=text
    )
    return np.array(response.data[0].embedding)

def enrich_tags(existing_tags, query_log):
    new_tags = set(existing_tags)
    for q in query_log.get("most_common_queries", []):
        for word in q.lower().split():
            if word not in new_tags and len(word) > 3:
                new_tags.add(word)
    return list(new_tags)

def map_columns_to_concepts(columns, global_aliases=None):
    unmapped = [col for col in columns if not global_aliases or col not in global_aliases]
    mapping = global_aliases.copy() if global_aliases else {}
    if unmapped:
        prompt = (
            "Map the following Excel column headers to standard business concepts. "
            "Return a JSON dictionary where each key is the original column name and each value is a standard concept "
            "(e.g., 'quantity', 'part_number', 'location'). If a header is junk or unrecognizable, map it to 'ignore'.\n\n"
            f"Columns: {unmapped}"
        )
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2
        )
        raw = response.choices[0].message.content.strip()
        try:
            new_mapping = json.loads(raw)
            mapping.update(new_mapping)
        except Exception:
            pass
    return mapping

def profile_excel(file_path):
    wb = openpyxl.load_workbook(file_path, read_only=True)
    profile = {}
    for sheet in wb.worksheets:
        sheet_profile = {}
        df = pd.DataFrame(sheet.values)
        df.columns = df.iloc[0]
        df = df[1:]
        for col in df.columns:
            col_data = df[col]
            sheet_profile[str(col)] = {
                "type": str(col_data.dtype),
                "null_pct": float(col_data.isnull().mean()),
                "min": col_data.min() if pd.api.types.is_numeric_dtype(col_data) else None,
                "max": col_data.max() if pd.api.types.is_numeric_dtype(col_data) else None,
                "sample_values": col_data.dropna().unique()[:5].tolist()
            }
        profile[sheet.title] = sheet_profile
    return profile

def scan_and_process_supabase(query_log=None):
    pass  # Supabase not used in this version
def extract_metadata_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    metadata = {}
    if prs.slides:
        first_slide = prs.slides[0]
        for shape in first_slide.shapes:
            if shape.has_text_frame:
                for line in shape.text_frame.text.split('\n'):
                    if ':' in line:
                        key, value = line.split(':', 1)
                        metadata[key.strip().lower()] = value.strip()
    return metadata

def extract_metadata_from_pdf(pdf_path):
    metadata = {}
    with pdfplumber.open(pdf_path) as pdf:
        lines = []
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                lines.extend(text.split('\n'))
            if len(lines) >= 10:
                break
        for line in lines[:10]:
            if ':' in line:
                key, value = line.split(':', 1)
                metadata[key.strip().lower()] = value.strip()
    return metadata

def extract_structural_metadata(raw_text, ext):
    if ext == ".docx":
        return {"section_headings": [line.strip() for line in raw_text.splitlines() if line.strip().endswith(":") or line.strip().istitle()]}
    elif ext == ".pdf":
        return {"section_headings": [line.strip() for line in raw_text.splitlines() if line.strip().endswith(":") or line.strip().istitle()]}
    elif ext == ".pptx":
        return {"slide_titles": [line.strip() for line in raw_text.splitlines() if len(line.strip()) > 0]}
    return {}

def generate_json(metadata, original_filename, profile_data=None, embedding=None, usage_stats=None, column_aliases=None):
    return {
        "title": metadata.get("title", ""),
        "author": metadata.get("author", ""),
        "date": metadata.get("date", str(datetime.today().date())),
        "category": metadata.get("category", ""),
        "tags": metadata.get("tags", []) if isinstance(metadata.get("tags"), list)
        else [tag.strip() for tag in metadata.get("tags", "").split(',') if tag.strip()] if "tags" in metadata else [],
        "summary": metadata.get("summary", ""),
        "source_file": original_filename,
        "profile_data": profile_data or {},
        "embedding": embedding.tolist() if embedding is not None else [],
        "usage_stats": usage_stats or {"times_used": 0, "last_used": None, "most_common_queries": []},
        "column_aliases": column_aliases or {}
    }

def write_sidecar_metadata(original_file_path, metadata_dict):
    folder = os.path.dirname(original_file_path)
    metadata_folder = os.path.join(folder, "_metadata")
    os.makedirs(metadata_folder, exist_ok=True)
    base = os.path.basename(original_file_path)
    meta_path = os.path.join(metadata_folder, f"{base}.json")
    with open(meta_path, "w", encoding="utf-8") as f:
        json.dump(metadata_dict, f, indent=2)
    return meta_path

# --- Streamlit UI trigger for local scan ---
st.header("LLM Metadata Sidecar Generator (Local)")
base_dir = st.text_input("Project root directory to scan:", value=os.getcwd())
if st.button("Scan and write metadata sidecars locally"):
    scan_and_process_local(base_dir)


