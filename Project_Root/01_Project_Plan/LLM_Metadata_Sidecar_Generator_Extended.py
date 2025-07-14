"""
LLM Metadata Sidecar Generator Script (Extended)
-------------------------------------------------
This script scans a project directory recursively for supported document types,
extracts structured metadata from the document itself (if present),
or falls back to OpenAI GPT to generate metadata.
It then saves a sidecar .json file into a `_metadata/` subfolder next to each file.

Supported: .docx, .xlsx, .pptx, .pdf
"""

import os
import json
from datetime import datetime
from docx import Document
import openpyxl
from pptx import Presentation
import pdfplumber
from openai import OpenAI
import PyPDF2
from pdfminer.high_level import extract_text
from pdf2image import convert_from_path
import pytesseract
import numpy as np
import pandas as pd

client = OpenAI(api_key="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")  # Replace with your real key

def get_embedding(text):
    response = client.embeddings.create(
        model="text-embedding-ada-002",
        input=text
    )
    return np.array(response.data[0].embedding)

def enrich_tags(existing_tags, query_log):
    # Add new tags from most common queries (simple example)
    new_tags = set(existing_tags)
    for q in query_log.get("most_common_queries", []):
        for word in q.lower().split():
            if word not in new_tags and len(word) > 3:
                new_tags.add(word)
    return list(new_tags)

def profile_excel(file_path):
    wb = openpyxl.load_workbook(file_path, read_only=True)
    profile = {}
    for sheet in wb.worksheets:
        sheet_profile = {}
        df = pd.DataFrame(sheet.values)
        df.columns = df.iloc[0]
        df = df[1:]
        for col in df.columns:
            col_data = df[col]
            sheet_profile[str(col)] = {
                "type": str(col_data.dtype),
                "null_pct": float(col_data.isnull().mean()),
                "min": col_data.min() if pd.api.types.is_numeric_dtype(col_data) else None,
                "max": col_data.max() if pd.api.types.is_numeric_dtype(col_data) else None,
                "sample_values": col_data.dropna().unique()[:5].tolist()
            }
        profile[sheet.title] = sheet_profile
    return profile

def extract_text_for_metadata(path):
    if path.endswith(".docx"):
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    elif path.endswith(".xlsx"):
        wb = openpyxl.load_workbook(path)
        text = []
        for sheet in wb.worksheets:
            text.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(max_row=5):  # Get first 5 rows for context
                row_text = [str(cell.value) for cell in row if cell.value]
                if row_text:
                    text.append(" | ".join(row_text))
        return "\n".join(text)
    elif path.endswith(".pptx"):
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text.append(shape.text)
        return "\n".join(text)
    elif path.endswith(".pdf"):
        text = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
        return "\n".join(text)
    return ""

def auto_generate_metadata(prompt_text):
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {
                "role": "system",
                "content": (
                    "You are a metadata assistant. Generate structured metadata for internal business documents.\n"
                    "Return only valid JSON with the following fields:\n"
                    "- title (short)\n"
                    "- category (1-2 words)\n"
                    "- tags (array of 3–7 keywords)\n"
                    "- summary (1 paragraph overview of the content; if the document is a spreadsheet, summarize what the data represents, including sheet and column names if possible)\n\n"
                    "Do not include markdown formatting, explanations, or anything other than valid JSON."
                )
            },
            {
                "role": "user",
                "content": f"Document content:\n{prompt_text[:3000]}\n\nReturn metadata JSON only."
            }
        ],
        temperature=0.2
    )

    raw = response.choices[0].message.content.strip()

    if raw.startswith("```json") or raw.startswith("```"):
        raw = raw.strip("`").replace("json", "").strip()

    try:
        metadata = json.loads(raw)
        metadata["title"] = metadata.get("title", "").strip()
        metadata["category"] = metadata.get("category", "").strip()
        metadata["tags"] = metadata.get("tags", []) if isinstance(metadata.get("tags", []), list) else []
        metadata["summary"] = metadata.get("summary", "").strip()
        return metadata
    except Exception as e:
        print("⚠️ GPT returned invalid JSON:", raw)
        raise e

def extract_metadata_from_docx(doc_path):
    doc = Document(doc_path)
    lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    metadata = {}
    for line in lines[:10]:
        if ':' in line:
            key, value = line.split(':', 1)
            metadata[key.strip().lower()] = value.strip()
    return metadata

def extract_metadata_from_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path, read_only=True)
    metadata = {}
    metadata["sheet_names"] = wb.sheetnames
    columns_by_sheet = {}
    for sheet in wb.worksheets:
        headers = [cell.value for cell in next(sheet.iter_rows(max_row=1))]
        columns_by_sheet[sheet.title] = [str(h) for h in headers if h]
    metadata["columns_by_sheet"] = columns_by_sheet
    all_columns = []
    for cols in columns_by_sheet.values():
        all_columns.extend(cols)
    metadata["columns"] = list(set(all_columns))
    ws = wb.active
    for row in ws.iter_rows(min_row=1, max_row=5, min_col=1, max_col=2):
        key = row[0].value
        value = row[1].value
        if key:
            metadata[str(key).strip().lower()] = str(value).strip() if value else ""
    return metadata

def extract_metadata_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    metadata = {}
    if prs.slides:
        first_slide = prs.slides[0]
        for shape in first_slide.shapes:
            if shape.has_text_frame:
                for line in shape.text_frame.text.split('\n'):
                    if ':' in line:
                        key, value = line.split(':', 1)
                        metadata[key.strip().lower()] = value.strip()
    return metadata

def extract_metadata_from_pdf(pdf_path):
    metadata = {}
    with pdfplumber.open(pdf_path) as pdf:
        lines = []
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                lines.extend(text.split('\n'))
            if len(lines) >= 10:
                break
        for line in lines[:10]:
            if ':' in line:
                key, value = line.split(':', 1)
                metadata[key.strip().lower()] = value.strip()
    return metadata

def generate_json(metadata, original_filename, profile_data=None, embedding=None, usage_stats=None):
    return {
        "title": metadata.get("title", ""),
        "author": metadata.get("author", ""),
        "date": metadata.get("date", str(datetime.today().date())),
        "category": metadata.get("category", ""),
        "tags": metadata.get("tags", []) if isinstance(metadata.get("tags"), list)
        else [tag.strip() for tag in metadata.get("tags", "").split(',') if tag.strip()] if "tags" in metadata else [],
        "summary": metadata.get("summary", ""),
        "source_file": original_filename,
        "profile_data": profile_data or {},
        "embedding": embedding.tolist() if embedding is not None else [],
        "usage_stats": usage_stats or {"times_used": 0, "last_used": None, "most_common_queries": []}
    }

def scan_and_process(base_dir, query_log=None):
    print(f"Scanning {base_dir}...")
    for root, dirs, files in os.walk(base_dir):
        for file in files:
            full_path = os.path.join(root, file)
            ext = os.path.splitext(file)[1].lower()
            base_name = os.path.splitext(file)[0]
            metadata = {}

            if ext == ".docx":
                metadata = extract_metadata_from_docx(full_path)
            elif ext == ".xlsx":
                metadata = extract_metadata_from_xlsx(full_path)
            elif ext == ".pptx":
                metadata = extract_metadata_from_pptx(full_path)
            elif ext == ".pdf":
                metadata = extract_metadata_from_pdf(full_path)
            else:
                continue

            # Always try to fill in missing fields (especially summary) with GPT
            raw_text = extract_text_for_metadata(full_path)
            gpt_metadata = auto_generate_metadata(raw_text[:3000] if raw_text else "")

            # Merge structured metadata and GPT metadata (GPT fills in missing fields)
            for key in ["title", "category", "tags", "summary"]:
                if not metadata.get(key):
                    metadata[key] = gpt_metadata.get(key, "")

            # --- Profile Data ---
            profile_data = {}
            if ext == ".xlsx":
                try:
                    profile_data = profile_excel(full_path)
                except Exception as e:
                    print(f"Profile data error: {e}")

            # --- Embedding ---
            summary_text = metadata.get("summary", "")
            embedding = get_embedding(summary_text) if summary_text else None

            # --- Usage Stats ---
            meta_dir = os.path.join(root, "_metadata")
            json_file_path = os.path.join(meta_dir, f"{base_name}.json")
            usage_stats = None
            if os.path.exists(json_file_path):
                try:
                    with open(json_file_path, "r") as f:
                        old = json.load(f)
                        usage_stats = old.get("usage_stats", None)
                        # Enrich tags based on query log if provided
                        if query_log and "tags" in metadata:
                            metadata["tags"] = enrich_tags(metadata["tags"], usage_stats or {})
                except Exception:
                    pass
            else:
                # If query_log provided, enrich tags
                if query_log and "tags" in metadata:
                    metadata["tags"] = enrich_tags(metadata["tags"], query_log)

            json_data = generate_json(metadata, file, profile_data, embedding, usage_stats)
            os.makedirs(meta_dir, exist_ok=True)
            with open(json_file_path, "w", encoding="utf-8") as f:
                json.dump(json_data, f, indent=2)
            print(f"✅ Metadata saved: {json_file_path}")

if __name__ == "__main__":
    # Optionally, pass a query_log dict for tag enrichment
    scan_and_process("/content/drive/MyDrive/Ethos LLM/Project_Root")