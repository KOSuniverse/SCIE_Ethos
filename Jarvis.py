import streamlit as st
import os
import json
from datetime import datetime
from openai import OpenAI
from docx import Document
import openpyxl
from pptx import Presentation
import pdfplumber
import pytesseract
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import joblib  # For loading prebuilt models
from io import BytesIO
from sharepoint_utils import (
    load_metadata, save_metadata,
    load_global_aliases, update_global_aliases,
    list_all_supported_files, download_file
)

client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

RAW_DATA_FOLDER = "04_Data/Raw_Data_Dumps/"
CLEANED_DATA_FOLDER = "04_Data/Cleaned_Engineering_Files/"
FINAL_ANALYSIS_FOLDER = "04_Data/Final_Analysis_Ready/"
MODELS_FOLDER = "04_Data/Models/"
METADATA_FOLDER = "01_Project_Plan/metadata/"
GLOBAL_ALIAS_FILE = "01_Project_Plan/global_column_aliases.json"

embedding_cache = {}

def get_embedding(text):
    try:
        response = client.embeddings.create(
            model="text-embedding-ada-002",
            input=text
        )
        return np.array(response.data[0].embedding)
    except Exception as e:
        st.error(f"Embedding error: {e}")
        return np.zeros(1536)

def cosine_similarity(a, b):
    if a is None or b is None or len(a) == 0 or len(b) == 0:
        return 0.0
    return float(np.dot(a, b) / (np.linalg.norm(a) * np.linalg.norm(b)))

def map_columns_to_concepts(columns, global_aliases=None):
    unmapped = [col for col in columns if not global_aliases or col not in global_aliases]
    mapping = global_aliases.copy() if global_aliases else {}
    if unmapped:
        prompt = (
            "Map the following Excel column headers to standard business concepts. "
            "Return a JSON dictionary where each key is the original column name and each value is a standard concept "
            "(e.g., 'quantity', 'part_number', 'location'). If a header is junk or unrecognizable, map it to 'ignore'.\n\n"
            f"Columns: {unmapped}"
        )
        try:
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.2
            )
            raw = response.choices[0].message.content.strip()
            new_mapping = json.loads(raw)
            mapping.update(new_mapping)
        except Exception as e:
            st.warning(f"Column mapping failed: {e}")
    return mapping

def extract_text_for_metadata(path, max_ocr_pages=5):
    ext = path.lower().split('.')[-1]
    try:
        file_stream = download_file(path)
        if ext == "docx":
            doc = Document(file_stream)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == "xlsx":
            wb = openpyxl.load_workbook(file_stream, read_only=True)
            text = []
            sheet_names = []
            columns_by_sheet = {}
            for sheet in wb.worksheets:
                sheet_names.append(sheet.title)
                text.append(f"Sheet: {sheet.title}")
                headers = [cell.value for cell in next(sheet.iter_rows(max_row=1))]
                columns_by_sheet[sheet.title] = [str(h) for h in headers if h]
                if headers:
                    text.append("Columns: " + " | ".join([str(h) for h in headers if h]))
            return "\n".join(text), sheet_names, columns_by_sheet
        elif ext == "pptx":
            prs = Presentation(file_stream)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text.append(shape.text)
            return "\n".join(text)
        elif ext == "pdf":
            text = []
            with pdfplumber.open(file_stream) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text.append(page_text)
            if any(t.strip() for t in text):
                return "\n".join(text)
            ocr_text = []
            with pdfplumber.open(file_stream) as pdf:
                for i, page in enumerate(pdf.pages):
                    if i >= max_ocr_pages:
                        break
                    st.write(f"OCR processing page {i+1}/{min(len(pdf.pages), max_ocr_pages)}")
                    image = page.to_image(resolution=300).original
                    page_ocr = pytesseract.image_to_string(image)
                    if page_ocr:
                        ocr_text.append(page_ocr)
            return "\n".join(ocr_text)
    except Exception as e:
        st.warning(f"Failed to extract text from {path}: {e}")
    return "", [], {}

def chunk_text(text, chunk_size=2000, overlap=200):
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = " ".join(words[i:i+chunk_size])
        chunks.append(chunk)
        i += chunk_size - overlap
    return chunks

# --- NEW: Robust LLM metadata and structure extraction ---
def generate_llm_metadata(text, file_type):
    prompt = (
        f"Analyze the following {file_type} content and provide:\n"
        "- A concise summary (2-3 sentences)\n"
        "- 3-7 relevant tags\n"
        "- A category (e.g., 'finance', 'engineering', 'presentation', etc.)\n\n"
        f"Content:\n{text[:4000]}"
    )
    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2
        )
        raw = response.choices[0].message.content.strip()
        try:
            meta = json.loads(raw)
        except Exception:
            meta = {}
            lines = raw.splitlines()
            for line in lines:
                if line.lower().startswith("summary:"):
                    meta["summary"] = line.split(":", 1)[1].strip()
                elif line.lower().startswith("tags:"):
                    meta["tags"] = [t.strip() for t in line.split(":", 1)[1].split(",")]
                elif line.lower().startswith("category:"):
                    meta["category"] = line.split(":", 1)[1].strip()
        return meta
    except Exception as e:
        st.warning(f"LLM metadata generation failed: {e}")
        return {}

def extract_structural_metadata(text, file_type):
    if file_type == "docx":
        return {"section_headings": [line.strip() for line in text.splitlines() if line.strip().endswith(":") or line.strip().istitle()]}
    elif file_type == "pdf":
        return {"section_headings": [line.strip() for line in text.splitlines() if line.strip().endswith(":") or line.strip().istitle()]}
    elif file_type == "pptx":
        return {"slide_titles": [line.strip() for line in text.splitlines() if len(line.strip()) > 0]}
    return {}

def get_file_last_modified(path):
    try:
        return os.path.getmtime(path)
    except Exception:
        return None

def verify_answer(answer, context, user_query):
    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a fact checker. Given a context and an answer, rate the factual accuracy (0-100) and list any unsupported claims."},
                {"role": "user", "content": f"Context:\n{context}\n\nAnswer:\n{answer}\n\nQuestion:\n{user_query}"}
            ],
            temperature=0.0
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"Verification failed: {e}"

def mine_for_root_causes(user_query, all_chunks, top_chunks):
    mining_prompt = (
        "You are a data mining assistant. Given a question and relevant context, "
        "search across all available data for patterns, correlations, or supporting evidence "
        "that could help explain the root cause of the issue or trend described in the question. "
        "Cite any files or data points that support your findings."
    )
    all_context = "\n\n".join([chunk for _, chunk in all_chunks])
    top_context = "\n\n".join([chunk for _, chunk in top_chunks])
    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": mining_prompt},
                {"role": "user", "content": f"Question:\n{user_query}\n\nRelevant context:\n{top_context}\n\nAll data:\n{all_context[:8000]}"}
            ],
            temperature=0.3
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"Root cause mining failed: {e}"

def predictive_modeling_prebuilt(user_query, top_chunks):
    model_files = [fp for fp in list_all_supported_files(MODELS_FOLDER) if fp.endswith('.pkl')]
    if not model_files:
        return "No prebuilt models found. Please build and upload a model first."
    try:
        model_path = model_files[0]
        model_stream = download_file(model_path)
        if not isinstance(model_stream, BytesIO):
            model_stream = BytesIO(model_stream.read())
        model = joblib.load(model_stream)
        return f"Prebuilt model '{os.path.basename(model_path)}' is available. Please implement feature extraction for real predictions."
    except Exception as e:
        return f"Could not load or run prebuilt model: {e}"

def predictive_modeling_guided(user_query, top_chunks):
    modeling_prompt = (
        "You are a predictive modeling assistant. Given a business question and relevant data, "
        "suggest an appropriate predictive model, describe the features to use, and provide Python code "
        "to train and use the model (e.g., scikit-learn, XGBoost). Explain the expected results and how this helps solve the business problem."
    )
    context = "\n\n".join([chunk for _, chunk in top_chunks])
    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": modeling_prompt},
                {"role": "user", "content": f"Question:\n{user_query}\n\nRelevant data:\n{context}"}
            ],
            temperature=0.3
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"Guided modeling failed: {e}"

def predictive_modeling_inference(user_query, top_chunks):
    inference_prompt = (
        "You are a business analyst. Given a question and relevant data, use statistical reasoning to simulate a prediction. "
        "Explain your reasoning and cite supporting data."
    )
    context = "\n\n".join([chunk for _, chunk in top_chunks])
    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": inference_prompt},
                {"role": "user", "content": f"Question:\n{user_query}\n\nRelevant data:\n{context}"}
            ],
            temperature=0.3
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"LLM-powered inference failed: {e}"

def excel_qa(file_path, user_query, column_aliases=None):
    import matplotlib.pyplot as plt
    import seaborn as sns
    import re

    try:
        file_stream = download_file(file_path)
        df = pd.read_excel(file_stream)
        auto_chart = any(word in user_query.lower() for word in ["trend", "compare", "distribution", "growth", "pattern", "chart", "plot", "visual"])
        prompt = (
            f"Column aliases for this file: {json.dumps(column_aliases or {})}\n"
            f"You are a data analyst working with the following Excel file.\n"
            f"Columns: {list(df.columns)}\n"
            f"Sample data:\n{df.head(5).to_string(index=False)}\n\n"
            f"User question: {user_query}\n\n"
            "Follow this reasoning chain:\n"
            "1. Identify the key data needed to answer the question.\n"
            "2. Retrieve or summarize the relevant data.\n"
            f"3. {'Generate a chart if it would help illustrate the answer (use matplotlib/seaborn and show it).' if auto_chart else 'Generate a chart if useful.'}\n"
            "4. Explain what the result or chart shows.\n"
            "5. Suggest 1–2 possible causes or business insights that could explain the observed pattern.\n\n"
            "Return only valid Python code that uses the provided 'df' DataFrame (do NOT reload or create new data). "
            "Assign any tabular result to a variable named 'result'.\n"
            "After the code block, provide your explanation and insights."
        )

        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a helpful data analyst."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.2
        )

        answer = response.choices[0].message.content.strip()
        code_match = re.search(r"```(?:python)?(.*?)```", answer, re.DOTALL)
        code = code_match.group(1).strip() if code_match else answer.strip()

        explanation = ""
        if code_match:
            explanation = answer[code_match.end():].strip()

        local_vars = {"df": df.copy(), "pd": pd, "plt": plt, "sns": sns}
        plt.clf()

        try:
            exec(code, {}, local_vars)
            if "result" in local_vars:
                st.write("✅ Result:")
                if isinstance(local_vars["result"], pd.DataFrame):
                    st.dataframe(local_vars["result"].astype(str))
                else:
                    st.write(local_vars["result"])
            st.pyplot(plt)
        except Exception as e:
            st.warning("⚠️ GPT-generated code did not execute successfully.")
            st.text(code)
            st.error(str(e))

        if explanation:
            st.markdown(f"**Explanation:** {explanation}")
    except Exception as e:
        st.error(f"Excel Q&A failed: {e}")

# --- Load global column aliases from SharePoint ---
global_aliases = load_global_aliases()

# --- Gather all files from Cleaned Engineering Files for Q&A and modeling ---
# List all supported files in the SharePoint document library, EXCLUDING raw data folder
all_files = [
    f["relative_path"]
    for f in list_all_supported_files()
    if not f["relative_path"].startswith(st.secrets["sharepoint"]["raw_data_folder"])
]

all_chunks = []
updated_global_aliases = global_aliases.copy()
progress_bar = st.progress(0, text="Indexing files and building metadata...")

for idx, file_path in enumerate(all_files):
    ext = file_path.lower().split('.')[-1]
    file_last_modified = get_file_last_modified(file_path)
    meta = load_metadata(file_path) or {"source_file": file_path}
    needs_index = True

    # Efficient re-indexing: skip if file hasn't changed since last_indexed
    if meta and meta.get("last_indexed") and file_last_modified:
        try:
            if file_last_modified <= datetime.fromisoformat(meta["last_indexed"]).timestamp():
                needs_index = False
        except Exception:
            pass

    if needs_index:
        try:
            if ext == "xlsx":
                text, sheet_names, columns_by_sheet = extract_text_for_metadata(file_path)
            else:
                text = extract_text_for_metadata(file_path)
                sheet_names, columns_by_sheet = [], {}
            if isinstance(text, tuple):
                text = text[0]
            if text.strip():
                meta["file_type"] = ext
                meta["file_size"] = os.path.getsize(file_path) if os.path.exists(file_path) else None
                meta["last_modified"] = datetime.fromtimestamp(file_last_modified).isoformat() if file_last_modified else None
                meta["last_indexed"] = datetime.now().isoformat()
                # Excel-specific
                if ext == "xlsx":
                    meta["sheet_names"] = sheet_names
                    meta["columns_by_sheet"] = columns_by_sheet
                    all_columns = []
                    for cols in columns_by_sheet.values():
                        all_columns.extend(cols)
                    meta["columns"] = list(set(all_columns))
                    column_aliases = map_columns_to_concepts(meta["columns"], updated_global_aliases)
                    meta["column_aliases"] = column_aliases
                    updated_global_aliases.update(column_aliases)
                # LLM-generated summary/tags/category for ALL files
                llm_meta = generate_llm_metadata(text, ext)
                meta.update(llm_meta)
                # Structural metadata for all files
                meta.update(extract_structural_metadata(text, ext))
                save_metadata(file_path, meta)
                for chunk in chunk_text(text):
                    all_chunks.append((meta, chunk))
                embedding = None
                if "embedding" in meta and meta["embedding"]:
                    embedding = np.array(meta["embedding"])
                elif meta.get("summary", ""):
                    try:
                        embedding = get_embedding(meta["summary"])
                    except Exception as e:
                        st.warning(f"Embedding failed for {file_path}: {e}")
                if embedding is not None:
                    embedding_cache[file_path] = embedding
        except Exception as e:
            st.warning(f"Failed to process {file_path}: {e}")
    else:
        # Use existing metadata and skip re-indexing
        if "summary" in meta:
            for chunk in chunk_text(meta["summary"]):
                all_chunks.append((meta, chunk))
        if "embedding" in meta and meta["embedding"]:
            embedding_cache[file_path] = np.array(meta["embedding"])
    progress_bar.progress((idx + 1) / len(all_files), text=f"Indexed {idx+1}/{len(all_files)} files")
progress_bar.empty()
update_global_aliases(updated_global_aliases)

if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "last_excel" not in st.session_state:
    st.session_state.last_excel = None
if "last_query" not in st.session_state:
    st.session_state.last_query = None
if "query_log" not in st.session_state:
    st.session_state.query_log = []

st.header("Ask a question about your knowledge base")

with st.form("question_form", clear_on_submit=False):
    user_query = st.text_input("Your question:", value=st.session_state.get("last_query", ""))
    submit = st.form_submit_button("Ask")

if submit and user_query.strip():
    st.info("Processing your question. This may take a moment...")
    scored_chunks = []
    query_embedding = get_embedding(user_query)
    for meta, chunk in all_chunks:
        meta_text = " ".join([
            str(meta.get("title", "")).lower(),
            str(meta.get("category", "")).lower(),
            " ".join(meta.get("tags", [])),
            str(meta.get("source_file", "")).lower(),
            str(meta.get("summary", "")).lower(),
            chunk.lower()
        ])
        keyword_score = sum(word in meta_text for word in user_query.lower().split())
        emb = embedding_cache.get(meta.get("source_file"))
        semantic_score = cosine_similarity(query_embedding, emb) if emb is not None else 0.0
        hybrid_score = 0.5 * keyword_score + 0.5 * semantic_score
        scored_chunks.append((hybrid_score, meta, chunk, keyword_score, semantic_score))

    scored_chunks.sort(reverse=True, key=lambda x: x[0])
    top_chunks = scored_chunks[:3] if scored_chunks else []

    context = "\n\n".join([chunk for _, _, chunk, _, _ in top_chunks])
    sources = [meta.get("source_file", "") for _, meta, _, _, _ in top_chunks]
    scores = [{"file": meta.get("source_file", ""), "score": score, "keyword": kw, "semantic": sem}
              for score, meta, _, kw, sem in top_chunks]

    system_prompt = (
        "You are a business analyst answering questions using internal documents (Excel, Word, PDF, PPTX). "
        "For each question, follow this reasoning chain:\n"
        "1. Identify the key data needed to answer the question.\n"
        "2. Retrieve or summarize the relevant information from the provided context.\n"
        "3. If the answer is quantitative and a chart would help, describe the chart (or provide code for matplotlib/seaborn if possible).\n"
        "4. Explain what the result or chart shows.\n"
        "5. Suggest possible causes or business insights that could explain the observed pattern, using supporting evidence from other loaded data if possible.\n"
        "6. Cite the source file(s) used for the answer.\n"
        "7. Rate your confidence in the answer (0-100) based on relevance scores and verification.\n"
        "Use only the provided context and cite sources."
    )

    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Context:\n{context}\n\nQuestion: {user_query}"}
            ],
            temperature=0.3
        )
        answer = response.choices[0].message.content.strip()
    except Exception as e:
        answer = f"LLM answer failed: {e}"

    verification = verify_answer(answer, context, user_query)
    root_cause_analysis = mine_for_root_causes(user_query, all_chunks, top_chunks)

    st.markdown("### Predictive Modeling Options")
    predictive_model_result = None
    if st.button("Run Prebuilt Model (if available)"):
        predictive_model_result = predictive_modeling_prebuilt(user_query, top_chunks)
        st.write(predictive_model_result)
    if st.button("Suggest Model & Code (LLM-guided)"):
        guided_model_code = predictive_modeling_guided(user_query, top_chunks)
        st.write(guided_model_code)
    if st.button("LLM-powered Inference (Simple Prediction)"):
        inference_result = predictive_modeling_inference(user_query, top_chunks)
        st.write(inference_result)

    st.markdown(f"**Answer:**\n\n{answer}")
    st.markdown("**Sources used:**")
    for src in sources:
        st.write(f"- {src}")
    st.markdown("**Relevance scores:**")
    for s in scores:
        st.write(f"{s['file']}: hybrid={s['score']:.2f}, keyword={s['keyword']}, semantic={s['semantic']:.2f}")
    st.markdown("**Verification:**")
    st.write(verification)
    st.markdown("**Root Cause/Data Mining Analysis:**")
    st.write(root_cause_analysis)

    top_meta = top_chunks[0][1] if top_chunks else None
    top_file = top_meta.get("source_file") if top_meta else None
    if top_file and top_file.lower().endswith('.xlsx'):
        st.markdown("---")
        st.markdown("**You can generate a chart from the top Excel file:**")
        if st.button("Show chart for top Excel file"):
            column_aliases = top_meta.get("column_aliases", {})
            excel_qa(top_file, user_query, column_aliases)

    st.session_state.last_query = user_query
    st.session_state.chat_history.append({"role": "user", "content": user_query})
    st.session_state.chat_history.append({"role": "assistant", "content": answer})
    st.session_state.query_log.append({
        "timestamp": datetime.now().isoformat(),
        "question": user_query,
        "files_used": sources,
        "scores": scores,
        "answer": answer,
        "verification": verification,
        "root_cause_analysis": root_cause_analysis,
        "predictive_model_result": predictive_model_result if predictive_model_result else "",
    })

else:
    st.write("Ask a question to get started.")

with st.expander("Show Query Log"):
    for entry in st.session_state.query_log:
        st.write(entry)

with st.expander("Show Chat History"):
    for msg in st.session_state.chat_history:
        st.write(f"{msg['role'].capitalize()}: {msg['content']}")