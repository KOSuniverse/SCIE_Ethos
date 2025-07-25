# --- Debug: Environment and Path Diagnostics ---
import getpass
import sys
import pathlib
import streamlit as st
import os
import json
import re
import warnings
import time
from datetime import datetime
from openai import OpenAI
from docx import Document
import openpyxl
from pptx import Presentation
import pdfplumber
import pytesseract
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from collections import defaultdict
import joblib  # For loading prebuilt models
from io import BytesIO
from difflib import SequenceMatcher
import traceback
import base64
import pickle
import io
import nltk
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

# --- Config ---
PROJECT_ROOT = r"C:\Users\dansk\OneDrive\Project_Root"
METADATA_FOLDER = os.path.join(PROJECT_ROOT, "01_Project_Plan", "_metadata")
GLOBAL_ALIAS_PATH = os.path.join(METADATA_FOLDER, "global_column_aliases.json")
LEARNED_ANSWERS_PATH = os.path.join(METADATA_FOLDER, "learned_answers.json")

SUPPORTED_EXTS = (".xlsx", ".docx", ".pdf", ".pptx")
OPENAI_EMBEDDING_MODEL = "text-embedding-ada-002"
OPENAI_CHAT_MODEL = "gpt-4o"

os.makedirs(METADATA_FOLDER, exist_ok=True)

client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

embedding_cache = {}

# --- File I/O Helpers ---

def find_all_supported_files(base_dir=PROJECT_ROOT):
    return [
        os.path.join(root, f)
        for root, dirs, files in os.walk(base_dir)
        for f in files if f.lower().endswith(SUPPORTED_EXTS)
    ]

def get_metadata_path(file_path):
    # Metadata file is always <document_folder>/_metadata/{basename}.json
    folder = os.path.dirname(file_path)
    metadata_folder = os.path.join(folder, "_metadata")
    os.makedirs(metadata_folder, exist_ok=True)
    base = os.path.basename(file_path)
    return os.path.join(metadata_folder, f"{base}.json")

def load_metadata(file_path):
    meta_path = get_metadata_path(file_path)
    if os.path.exists(meta_path):
        with open(meta_path, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}

def save_metadata(file_path, meta):
    meta_path = get_metadata_path(file_path)
    os.makedirs(os.path.dirname(meta_path), exist_ok=True)  # Ensure folder exists
    try:
        with open(meta_path, "w", encoding="utf-8") as f:
            json.dump(meta, f, indent=2)
        st.info(f"Metadata saved to {meta_path}")
    except Exception as e:
        st.error(f"Failed to save metadata to {meta_path}: {e}")

def load_global_aliases():
    if os.path.exists(GLOBAL_ALIAS_PATH):
        with open(GLOBAL_ALIAS_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}

def save_global_aliases(aliases):
    with open(GLOBAL_ALIAS_PATH, "w", encoding="utf-8") as f:
        json.dump(aliases, f, indent=2)

def load_learned_answers():
    if os.path.exists(LEARNED_ANSWERS_PATH):
        with open(LEARNED_ANSWERS_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}

def save_learned_answers(answers):
    with open(LEARNED_ANSWERS_PATH, "w", encoding="utf-8") as f:
        json.dump(answers, f, indent=2)

def download_file(file_path):
    # Local usage: just return file path
    return file_path

# --- Embeddings & Scoring ---

def openai_with_retry(call_fn, max_retries=3, delay=2):
    for attempt in range(max_retries):
        try:
            return call_fn()
        except Exception as e:
            if attempt < max_retries - 1:
                time.sleep(delay)
            else:
                raise e

def validate_embedding(vec, expected_dim=1536):
    if vec is None or not isinstance(vec, np.ndarray) or vec.shape[0] != expected_dim:
        return np.zeros(expected_dim)
    return vec

def get_embedding(text):
    try:
        response = openai_with_retry(
            lambda: client.embeddings.create(
                model=OPENAI_EMBEDDING_MODEL,
                input=text
            )
        )
        return validate_embedding(np.array(response.data[0].embedding))
    except Exception as e:
        st.error(f"Embedding error: {e}")
        return np.zeros(1536)

def cosine_similarity(a, b):
    if a is None or b is None or len(a) == 0 or len(b) == 0:
        return 0.0
    return float(np.dot(a, b) / (np.linalg.norm(a) * np.linalg.norm(b)))

# --- Metadata and Column Alias Mapping ---

def detect_alias_conflicts(mapping):
    """
    Detect conflicts where multiple columns map to the same alias.
    Returns a dict of alias -> [columns] for conflicts (excluding 'ignore').
    """
    reverse = defaultdict(list)
    for col, alias in mapping.items():
        reverse[alias].append(col)
    conflicts = {
        alias: cols for alias, cols in reverse.items()
        if len(cols) > 1 and alias != "ignore"
    }
    return conflicts

def fuzzy_match_score(word, text, threshold=0.85):
    """
    Calculate fuzzy match score for keyword matching.
    Returns True if any token in text matches word above threshold.
    """
    from difflib import SequenceMatcher
    return any(
        SequenceMatcher(None, word, token).ratio() >= threshold
        for token in text.split()
    )

def map_columns_to_concepts(columns, global_aliases=None, preview=True):
    """
    Map Excel column headers to standardized business concepts using LLM.
    Args:
        columns: List of column names to map
        global_aliases: Existing global alias mappings
        preview: Whether to show preview/editing interface
    Returns:
        Dictionary mapping column names to concepts
    """
    if not columns:
        return {}
    unmapped = [col for col in columns if not global_aliases or col not in global_aliases]
    mapping = global_aliases.copy() if global_aliases else {}
    new_mapping = {}
    if unmapped:
        prompt = (
            "You are a business data normalization assistant. Given spreadsheet column headers, map each to the best business alias from these choices: "
            "['part_number', 'description', 'category', 'location', 'unit_of_measure', 'quantity', 'value', 'price', 'date', 'period', 'order_number', 'customer', 'supplier', 'account_number', 'manager', 'owner', 'line_number', 'status', 'activity', 'reference_number', 'variance', 'cost_center', 'invoice', 'receipt', 'burden', 'labor_cost', 'labor_hours', 'forecast', 'lead_time', 'inventory_turns', 'sales', 'project', 'budget', 'actual', 'plan', 'balance', 'period_label', 'cost', 'currency', 'company']. "
            "If none apply, return 'ignore'. Only return a JSON dictionary, no explanation.\n\n"
            "Example output:\n{\n  'QTY': 'quantity',\n  'part no': 'part_number',\n  'xyz123': 'ignore'\n}\n\n"
            f"Columns to map: {unmapped}\n"
            "Only return valid JSON. Do not include explanations, comments, or markdown formatting."
        )
        try:
            response = openai_with_retry(
                lambda: client.chat.completions.create(
                    model=OPENAI_CHAT_MODEL,
                    messages=[{"role": "user", "content": prompt}],
                    temperature=0.2
                )
            )
            raw = response.choices[0].message.content.strip()
            # Remove markdown wrapper if present
            if raw.startswith("```json") or raw.startswith("```"):
                raw = raw.strip("`").replace("json", "").strip()
            # Try to parse JSON
            try:
                new_mapping = json.loads(raw)
            except Exception:
                # Fallback: try to parse as key-value pairs
                new_mapping = {}
                for line in raw.splitlines():
                    if ":" in line:
                        k, v = line.split(":", 1)
                        new_mapping[k.strip().strip('"')] = v.strip().strip('"').rstrip(',')
                if not new_mapping:
                    st.warning("Could not parse column mapping from LLM output.")
            mapping.update(new_mapping)
        except Exception as e:
            st.warning(f"Column mapping failed: {e}")
    # Detect alias conflicts
    conflicts = detect_alias_conflicts(mapping)
    if conflicts:
        st.warning(f"Alias conflicts detected: {conflicts}")
    # Preview/audit in Streamlit
    if preview and new_mapping:
        st.write("Column mapping preview (edit if needed):")
        editable_json = st.text_area(
            "Edit mapping as JSON if needed:",
            value=json.dumps(mapping, indent=2),
            key=f"column_mapping_preview_{hash(str(columns))}"
        )
        try:
            mapping = json.loads(editable_json)
        except Exception:
            st.error("Invalid JSON in edited mapping. Using previous mapping.")
    return mapping

def safe_llm_input(text, summary=None, max_len=4000):
    if isinstance(text, str) and len(text) > max_len:
        return summary if summary else text[:max_len]
    return text

def extract_text_for_metadata(path, max_ocr_pages=5):
    ext = path.lower().split('.')[-1]
    try:
        if ext == "docx":
            doc = Document(path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return text, [], {}
        elif ext == "xlsx":
            wb = openpyxl.load_workbook(path, read_only=True)
            text = []
            sheet_names = []
            columns_by_sheet = {}
            for sheet in wb.worksheets:
                sheet_names.append(sheet.title)
                text.append(f"Sheet: {sheet.title}")
                headers = [cell.value for cell in next(sheet.iter_rows(max_row=1))]
                columns_by_sheet[sheet.title] = [str(h) for h in headers if h]
                if headers:
                    text.append("Columns: " + " | ".join(columns_by_sheet[sheet.title]))
            return "\n".join(text), sheet_names, columns_by_sheet
        elif ext == "pptx":
            prs = Presentation(path)
            text = "\n".join(
                shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame
            )
            return text, [], {}
        elif ext == "pdf":
            cache_dir = os.path.join(os.path.dirname(path), "_ocr_cache")
            os.makedirs(cache_dir, exist_ok=True)
            cache_file = os.path.join(cache_dir, os.path.splitext(os.path.basename(path))[0] + ".txt")
            if os.path.exists(cache_file):
                with open(cache_file, "r", encoding="utf-8") as f:
                    return f.read(), [], {}
            with pdfplumber.open(path) as pdf:
                extracted_text = [page.extract_text() for page in pdf.pages if page.extract_text()]
                if any(extracted_text):
                    result = "\n".join(extracted_text)
                    with open(cache_file, "w", encoding="utf-8") as f:
                        f.write(result)
                    return result, [], {}
                # OCR fallback
                ocr_text = []
                for i, page in enumerate(pdf.pages[:max_ocr_pages]):
                    st.write(f"OCR processing page {i+1}/{min(len(pdf.pages), max_ocr_pages)}")
                    image = page.to_image(resolution=300).original
                    ocr_page = pytesseract.image_to_string(image)
                    if ocr_page:
                        ocr_text.append(ocr_page)
                result = "\n".join(ocr_text)
                with open(cache_file, "w", encoding="utf-8") as f:
                    f.write(result)
                return result, [], {}
    except Exception as e:
        st.warning(f"Failed to extract text from {path}: {e}")
    return "", [], {}

def chunk_text(text, chunk_size=2000, overlap=200, max_chars=4000):
    sentences = nltk.sent_tokenize(text)
    chunks = []
    current_chunk = []
    current_len = 0
    start_idx = 0

    for sent in sentences:
        sent = sent.strip()
        if not sent:
            continue
        if current_len + len(sent) > max_chars or len(current_chunk) >= chunk_size:
            chunk_text = " ".join(current_chunk).strip()
            end_idx = start_idx + len(chunk_text)
            chunks.append((chunk_text, start_idx, end_idx))
            start_idx = end_idx
            # Overlap logic
            if overlap > 0 and len(current_chunk) > overlap:
                current_chunk = current_chunk[-overlap:]
                current_len = sum(len(s) for s in current_chunk)
            else:
                current_chunk = []
                current_len = 0
        current_chunk.append(sent)
        current_len += len(sent)
    if current_chunk:
        chunk_text = " ".join(current_chunk).strip()
        end_idx = start_idx + len(chunk_text)
        chunks.append((chunk_text, start_idx, end_idx))
    return chunks

def generate_llm_metadata(text, file_type):
    prompt = f"""
You are an expert document classification system trained to extract structured metadata from internal business documents.

Your task is to analyze the following file content and return metadata that will be used to:
- Select relevant documents in response to business questions
- Embed file context for AI search
- Categorize content for analytics and modeling

📄 File Type: {file_type}

📝 Analyze the content below (first 4000 characters):
{text[:4000]}

🎯 Return structured JSON metadata in the format:
{{
  "title": "Short descriptive title (10 words max)",
  "summary": "Detailed overview of the document content, including its business purpose, key fields, sheet/section names, and any key metrics or concepts. Use 4–6 sentences or more if needed.",
  "tags": [
    "specific keyword or metric",
    "synonym or variant if applicable",
    "e.g., 'par levels', 'usage', 'utilization'",
    "limit 7–10 tags",
    "avoid generic adjectives"
  ],
  "category": "Primary topic or department (e.g., 'supply_chain', 'finance', 'compliance')"
}}

📌 Rules:
- Do not guess or invent details not supported by content
- Tags must reflect exact terms OR common business synonyms
- Summary must be informative for AI or human previewers
- Output only valid JSON — no markdown, explanation, or comments
"""
    try:
        response = client.chat.completions.create(
            model=OPENAI_CHAT_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2
        )
        raw = response.choices[0].message.content.strip()
        if raw.startswith("```json") or raw.startswith("```"):
            raw = raw.strip("`").replace("json", "").strip()
        return json.loads(raw)
    except Exception as e:
        st.warning(f"LLM metadata generation failed: {e}")
        return {}

# All other advanced logic (root cause, predictive, Q&A, UI) — unchanged from your Jarvis(1) version,
# but all file/metadata calls now use the above helpers.

# ... (Continued: the rest of the code is very long! See next message) ...
# --- Predictive Modeling & Root Cause ---

def verify_answer(answer, context, user_query):
    try:
        response = client.chat.completions.create(
            model=OPENAI_CHAT_MODEL,
            messages=[
                {"role": "system", "content": "You are a fact checker. Given a context and an answer, rate the factual accuracy (0-100) and list any unsupported claims."},
                {"role": "user", "content": f"Context:\n{context}\n\nAnswer:\n{answer}\n\nQuestion:\n{user_query}"}
            ],
            temperature=0.0
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"Verification failed: {e}"

def mine_for_root_causes(user_query, all_chunks, top_chunks):
    mining_prompt = """
You are a senior data analyst trained to identify operational root causes across business documents such as Excel files, reports, audits, or presentations.

Your task is to analyze the provided context in response to a user’s question. Using supporting content from the data, identify possible root causes, key drivers, or anomalies. These may relate to inventory problems, demand changes, usage patterns, compliance issues, or other operational signals.

Return your analysis in this format (JSON only):
{
  "root_cause_summary": "Concise paragraph explaining the issue",
  "supporting_evidence": [
    {"file": "filename.xlsx", "reason": "Column 'usage' drops in March while 'inventory' spikes"},
    {"file": "audit_report.pdf", "reason": "Notes a stock adjustment not reflected in system"}
  ],
  "confidence_score": 0–100
}
Do not include markdown, comments, or speculative guesses without data support.
"""
    all_context = "\n\n".join([chunk for _, chunk in all_chunks])
    top_context = "\n\n".join([chunk for _, _, chunk, _, _ in top_chunks])
    try:
        response = client.chat.completions.create(
            model=OPENAI_CHAT_MODEL,
            messages=[
                {"role": "system", "content": mining_prompt},
                {"role": "user", "content": f"Question:\n{user_query}\n\nRelevant context:\n{top_context}\n\nAll data:\n{all_context[:8000]}"}
            ],
            temperature=0.3
        )
        raw = response.choices[0].message.content.strip()
        if raw.startswith("```json") or raw.startswith("```"):
            raw = raw.strip("`").replace("json", "").strip()
        try:
            return json.loads(raw)
        except Exception:
            return raw
    except Exception as e:
        return f"Root cause mining failed: {e}"

# --- Predictive Modeling ---

def predictive_modeling_prebuilt(user_query, top_chunks):
    model_files = [os.path.join(PROJECT_ROOT, "04_Data", "Models", f) 
                   for f in os.listdir(os.path.join(PROJECT_ROOT, "04_Data", "Models")) if f.endswith('.pkl')]
    if not model_files:
        return "No prebuilt models found. Please build and upload a model first."
    try:
        model_path = model_files[0]
        with open(model_path, "rb") as f:
            model = joblib.load(f)
        return f"Prebuilt model '{os.path.basename(model_path)}' is available. Please implement feature extraction for real predictions."
    except Exception as e:
        return f"Could not load or run prebuilt model: {e}"

def get_model_prompt_response(user_query, context_df_sample):
    modeling_prompt = f"""
You are a senior machine learning engineer. You are given a business question and structured tabular data from Excel files.

Your task is to:
1. Identify the prediction goal (e.g., regression, classification, forecasting).
2. Choose the target variable and relevant features from the dataset.
3. Generate a robust Python model training pipeline using scikit-learn or XGBoost.
4. Include data cleaning, encoding, train/test split, training, evaluation, and prediction on test data.

Use the DataFrame variable `df` as your data.

🧾 Sample of the data:
{context_df_sample.to_string(index=False)}

📄 User Question:
{user_query}

Return valid JSON:
{{
  "model_type": "regression | classification | forecasting | anomaly_detection",
  "target_variable": "string",
  "features": ["col1", "col2", ...],
  "model_code": "FULL Python code as string that uses df",
  "description": "Explanation of the model design"
}}
Do not include markdown or comments. JSON only.
"""
    try:
        response = client.chat.completions.create(
            model=OPENAI_CHAT_MODEL,
            messages=[{"role": "user", "content": modeling_prompt}],
            temperature=0.3
        )
        raw = response.choices[0].message.content.strip()
        if raw.startswith("```json"):
            raw = raw.strip("`").replace("json", "").strip()
        return json.loads(raw)
    except Exception as e:
        st.warning("❌ GPT model generation failed.")
        st.text(traceback.format_exc())
        return None

def run_model_code(model_code, df):
    local_vars = {"df": df.copy()}
    try:
        exec(model_code, {}, local_vars)
        return local_vars
    except Exception as e:
        st.error("❌ Failed to execute generated model code.")
        st.text(model_code)
        st.text(traceback.format_exc())
        return None

def save_model(local_vars, model_name="auto_model.pkl"):
    model = local_vars.get("model", None)
    if model:
        model_bytes = pickle.dumps(model)
        models_dir = os.path.join(PROJECT_ROOT, "04_Data", "Models")
        os.makedirs(models_dir, exist_ok=True)
        path = os.path.join(models_dir, model_name)
        with open(path, "wb") as f:
            f.write(model_bytes)
        st.success(f"✅ Model saved to {path}")
    else:
        st.warning("⚠️ No model object found to save.")

def build_and_run_model(user_query, top_chunks, target_column=None):
    try:
        sample_df = None
        for meta, chunk in top_chunks:
            file = meta.get("source_file", "")
            if file.endswith(".xlsx"):
                file_stream = download_file(file)
                sample_df = pd.read_excel(file_stream)
                break
        if sample_df is None:
            return "⚠️ No usable Excel data found for modeling."

        result = get_model_prompt_response(user_query, sample_df.head(5))
        if not result:
            return "❌ Model generation failed."

        model_code = result.get("model_code", "")
        st.subheader("🧠 GPT-Generated Model")
        st.code(model_code, language="python")

        local_vars = run_model_code(model_code, sample_df)
        if not local_vars:
            return "⚠️ Code execution failed."

        save_model(local_vars)

        if "y_pred" in local_vars:
            st.subheader("📈 Model Predictions")
            st.dataframe(pd.DataFrame(local_vars["y_pred"], columns=["Prediction"]))
    except Exception as e:
        st.error(f"Model building/running failed: {e}")

# --- Excel Q&A and Charting ---


# --- Enhanced Excel Q&A Workflow ---
def remap_columns(df, column_aliases):
    mapping = {}
    if column_aliases:
        for col in df.columns:
            std = column_aliases.get(col)
            if std and std != "ignore":
                mapping[col] = std
    return df.rename(columns=mapping)

def llm_direct_answer(df, user_query, n=5):
    """
    Sends prompt to LLM, gets direct answer as a table (not code or explanation).
    """
    region_col = None
    for col in ["region", "location", "site", "warehouse"]:
        if col in df.columns:
            region_col = col
            break
    table_head = df.head(10).to_string(index=False)
    prompt = (
        f"You are a data analyst working with the following data from an Excel file.\n"
        f"Columns: {list(df.columns)}\n"
        f"Sample data:\n{table_head}\n\n"
        f"User question: {user_query}\n\n"
        f"Return ONLY the direct answer as a markdown table. "
        f"If a region column is present, break out the result by region. "
        f"Limit to the top {n} results if appropriate. "
        f"Do NOT include code, explanations, or commentary."
    )
    response = client.chat.completions.create(
        model=OPENAI_CHAT_MODEL,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.2
    )
    return response.choices[0].message.content.strip()

def llm_button_insight(df, user_query, insight_type, context_table=None):
    base = (
        f"Columns: {list(df.columns)}\n"
        f"User question: {user_query}\n\n"
    )
    if context_table:
        base += f"Table data:\n{context_table}\n\n"
    prompt_map = {
        "chart": "Suggest the most effective chart to visualize this result. Do NOT provide code—just describe the chart type, axes, and what it will show.",
        "explain": "Explain what this result means for the business. Focus on the practical significance of the findings.",
        "causes": "Suggest 1–2 possible business causes or drivers for the observed pattern in this result.",
        "root_cause": "Analyze the data and result to identify possible root causes, anomalies, or business drivers. Provide a concise summary."
    }
    prompt = base + prompt_map[insight_type]
    response = client.chat.completions.create(
        model=OPENAI_CHAT_MODEL,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.2
    )
    return response.choices[0].message.content.strip()

def excel_qa(file_path, user_query, column_aliases=None, n=5):
    # --- 1. Load and alias DataFrame
    df = pd.read_excel(file_path)
    df = remap_columns(df, column_aliases)

    # --- 2. Get direct answer as markdown table from LLM
    direct_answer_md = llm_direct_answer(df, user_query, n=n)
    st.markdown("### Answer")
    st.markdown(direct_answer_md)

    # Store table context for followup analysis
    st.session_state["last_answer_table"] = direct_answer_md
    st.session_state["last_qa_df"] = df
    st.session_state["last_qa_query"] = user_query

    # --- 3. Layered insights via buttons
    st.markdown("---")
    st.markdown("#### Drill Down:")

    if st.button("📊 Show Chart"):
        chart_desc = llm_button_insight(df, user_query, "chart", context_table=direct_answer_md)
        st.markdown(f"**Chart Recommendation:**\n{chart_desc}")

    if st.button("💡 Explain Result"):
        explanation = llm_button_insight(df, user_query, "explain", context_table=direct_answer_md)
        st.markdown(f"**Explanation:**\n{explanation}")

    if st.button("🧐 Suggest Causes"):
        causes = llm_button_insight(df, user_query, "causes", context_table=direct_answer_md)
        st.markdown(f"**Suggested Causes:**\n{causes}")

    if st.button("🔍 Root Cause/Data Mining"):
        root_cause = llm_button_insight(df, user_query, "root_cause", context_table=direct_answer_md)
        st.markdown(f"**Root Cause Analysis:**\n{root_cause}")

# --- Streamlit UI and App Main Loop ---

if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "last_excel" not in st.session_state:
    st.session_state.last_excel = None
if "last_query" not in st.session_state:
    st.session_state.last_query = None
if "query_log" not in st.session_state:
    st.session_state.query_log = []


# --- Robust Multi-Alias Mapping ---
def map_columns_to_multi_aliases(columns, concept_alias_dict=None, preview=True):
    """
    Map Excel column headers to standardized business concepts, supporting multiple aliases per concept.
    Args:
        columns: List of column names to map
        concept_alias_dict: Existing concept->aliases dict
        preview: Whether to show preview/editing interface
    Returns:
        Dictionary mapping concept to list of aliases
    """
    if not columns:
        return {}
    # Build reverse lookup: alias->concept
    concept_alias_dict = concept_alias_dict or {}
    # Flatten all known aliases
    known_aliases = set()
    for alias_list in concept_alias_dict.values():
        if isinstance(alias_list, list):
            known_aliases.update(alias_list)
        elif isinstance(alias_list, str):
            known_aliases.add(alias_list)
    # Find new columns not already mapped
    unmapped = [col for col in columns if col not in known_aliases]
    concept_to_aliases = {k: v if isinstance(v, list) else [v] for k, v in concept_alias_dict.items()}
    new_aliases = defaultdict(list)
    if unmapped:
        # Use LLM to map each column to a concept and collect robust aliases
        prompt = (
            "You are a business data normalization assistant. Given spreadsheet column headers, map each to the best business concept from these choices: "
            "['part_number', 'description', 'category', 'location', 'unit_of_measure', 'quantity', 'value', 'price', 'date', 'period', 'order_number', 'customer', 'supplier', 'account_number', 'manager', 'owner', 'line_number', 'status', 'activity', 'reference_number', 'variance', 'cost_center', 'invoice', 'receipt', 'burden', 'labor_cost', 'labor_hours', 'forecast', 'lead_time', 'inventory_turns', 'sales', 'project', 'budget', 'actual', 'plan', 'balance', 'period_label', 'cost', 'currency', 'company']. "
            "For each concept, return a list of all possible header variants (synonyms, abbreviations, alternate spellings, etc). If none apply, return 'ignore'. Only return a JSON dictionary mapping concept to list of aliases."
            "Example output:\n{\n  'part_number': ['part_no', 'part number', 'pn', 'item_no', ...],\n  'description': ['description', 'desc', 'details'],\n  ...\n}\n"
            f"Columns to map: {unmapped}\n"
            "Only return valid JSON. Do not include explanations, comments, or markdown formatting."
        )
        try:
            response = openai_with_retry(
                lambda: client.chat.completions.create(
                    model=OPENAI_CHAT_MODEL,
                    messages=[{"role": "user", "content": prompt}],
                    temperature=0.2
                )
            )
            raw = response.choices[0].message.content.strip()
            if raw.startswith("```json") or raw.startswith("```"):
                raw = raw.strip("`").replace("json", "").strip()
            try:
                new_aliases = json.loads(raw)
            except Exception:
                # Fallback: try to parse as key-value pairs
                new_aliases = defaultdict(list)
                for line in raw.splitlines():
                    if ":" in line:
                        k, v = line.split(":", 1)
                        k = k.strip().strip('"')
                        v = v.strip().strip('"').rstrip(',')
                        if v.startswith("[") and v.endswith("]"):
                            try:
                                v_list = json.loads(v)
                                new_aliases[k] = v_list
                            except Exception:
                                new_aliases[k] = [v]
                        else:
                            new_aliases[k] = [v]
                if not new_aliases:
                    st.warning("Could not parse multi-alias mapping from LLM output.")
            # Merge new aliases into concept_to_aliases
            for concept, aliases in new_aliases.items():
                if concept in concept_to_aliases:
                    # Merge, avoid duplicates
                    concept_to_aliases[concept] = list(set(concept_to_aliases[concept]) | set(aliases))
                else:
                    concept_to_aliases[concept] = aliases
        except Exception as e:
            st.warning(f"Multi-alias mapping failed: {e}")
    # Preview/audit in Streamlit
    if preview and new_aliases:
        st.write("Multi-alias mapping preview (edit if needed):")
        editable_json = st.text_area(
            "Edit mapping as JSON if needed:",
            value=json.dumps(concept_to_aliases, indent=2),
            key=f"multi_alias_mapping_preview_{hash(str(columns))}"
        )
        try:
            concept_to_aliases = json.loads(editable_json)
        except Exception:
            st.error("Invalid JSON in edited mapping. Using previous mapping.")
    return concept_to_aliases

# --- Reindex Logic ---
def reindex_all_files():
    """
    Full reindex: scan all files, remap columns with multi-alias support, refresh metadata, update alias file.
    """
    all_files = find_all_supported_files(PROJECT_ROOT)
    all_chunks = []
    global_aliases = load_global_aliases()
    updated_global_aliases = global_aliases.copy() if global_aliases else {}
    progress_bar = st.progress(0, text="Reindexing files and building metadata...")
    for idx, file_path in enumerate(all_files):
        ext = file_path.lower().split('.')[-1]
        meta = load_metadata(file_path) or {"source_file": file_path}
        try:
            if ext == "xlsx":
                text, sheet_names, columns_by_sheet = extract_text_for_metadata(file_path)
            else:
                text = extract_text_for_metadata(file_path)[0]
                sheet_names, columns_by_sheet = [], {}
            if text.strip():
                meta["file_type"] = ext
                meta["file_size"] = os.path.getsize(file_path) if os.path.exists(file_path) else None
                meta["last_modified"] = datetime.fromtimestamp(os.path.getmtime(file_path)).isoformat() if os.path.exists(file_path) else None
                meta["last_indexed"] = datetime.now().isoformat()
                if ext == "xlsx":
                    meta["sheet_names"] = sheet_names
                    meta["columns_by_sheet"] = columns_by_sheet
                    all_columns = []
                    for cols in columns_by_sheet.values():
                        all_columns.extend(cols)
                    meta["columns"] = list(set(all_columns))
                    # Multi-alias mapping
                    column_aliases = map_columns_to_multi_aliases(meta["columns"], updated_global_aliases, preview=False)
                    meta["column_aliases"] = column_aliases
                    updated_global_aliases = column_aliases
                llm_meta = generate_llm_metadata(text, ext)
                meta.update(llm_meta)
                save_metadata(file_path, meta)
                for chunk in chunk_text(text):
                    all_chunks.append((meta, chunk[0]))
                embedding = None
                if "embedding" in meta and meta["embedding"]:
                    embedding = np.array(meta["embedding"])
                elif meta.get("summary", ""):
                    try:
                        embedding = get_embedding(meta["summary"])
                    except Exception as e:
                        st.warning(f"Embedding failed for {file_path}: {e}")
                if embedding is not None:
                    embedding_cache[file_path] = embedding
        except Exception as e:
            st.warning(f"Failed to process {file_path}: {e}")
        progress_bar.progress((idx + 1) / len(all_files), text=f"Reindexed {idx+1}/{len(all_files)} files")
    progress_bar.empty()
    save_global_aliases(updated_global_aliases)
    st.success("Reindex complete! All mappings and metadata have been refreshed.")

# --- UI: Title + Reindex Button ---
col1, col2 = st.columns([4,1])
with col1:
    st.header("Ask a question about your knowledge base")
with col2:
    if st.button("🔄 Reindex All Files"):
        reindex_all_files()

# --- Main App Logic (unchanged, but use multi-alias mapping for new files) ---
all_files = find_all_supported_files(PROJECT_ROOT)
all_chunks = []
global_aliases = load_global_aliases()
updated_global_aliases = global_aliases.copy() if global_aliases else {}
progress_bar = st.progress(0, text="Indexing files and building metadata...")

new_columns_mapped = False
for idx, file_path in enumerate(all_files):
    ext = file_path.lower().split('.')[-1]
    meta = load_metadata(file_path) or {"source_file": file_path}
    needs_index = True
    file_last_modified = os.path.getmtime(file_path) if os.path.exists(file_path) else None
    if meta and meta.get("last_indexed") and file_last_modified:
        try:
            if file_last_modified <= datetime.fromisoformat(meta["last_indexed"]).timestamp():
                needs_index = False
        except Exception:
            pass
    if needs_index:
        try:
            if ext == "xlsx":
                text, sheet_names, columns_by_sheet = extract_text_for_metadata(file_path)
            else:
                text = extract_text_for_metadata(file_path)[0]
                sheet_names, columns_by_sheet = [], {}
            if text.strip():
                meta["file_type"] = ext
                meta["file_size"] = os.path.getsize(file_path) if os.path.exists(file_path) else None
                meta["last_modified"] = datetime.fromtimestamp(file_last_modified).isoformat() if file_last_modified else None
                meta["last_indexed"] = datetime.now().isoformat()
                if ext == "xlsx":
                    meta["sheet_names"] = sheet_names
                    meta["columns_by_sheet"] = columns_by_sheet
                    all_columns = []
                    for cols in columns_by_sheet.values():
                        all_columns.extend(cols)
                    meta["columns"] = list(set(all_columns))
                    prev_aliases = set(updated_global_aliases.keys())
                    column_aliases = map_columns_to_multi_aliases(meta["columns"], updated_global_aliases, preview=True)
                    meta["column_aliases"] = column_aliases
                    updated_global_aliases = column_aliases
                    if set(meta["columns"]) - prev_aliases:
                        new_columns_mapped = True
                llm_meta = generate_llm_metadata(text, ext)
                meta.update(llm_meta)
                save_metadata(file_path, meta)
                for chunk in chunk_text(text):
                    all_chunks.append((meta, chunk[0]))
                embedding = None
                if "embedding" in meta and meta["embedding"]:
                    embedding = np.array(meta["embedding"])
                elif meta.get("summary", ""):
                    try:
                        embedding = get_embedding(meta["summary"])
                    except Exception as e:
                        st.warning(f"Embedding failed for {file_path}: {e}")
                if embedding is not None:
                    embedding_cache[file_path] = embedding
        except Exception as e:
            st.warning(f"Failed to process {file_path}: {e}")
    else:
        if "summary" in meta:
            for chunk in chunk_text(meta["summary"]):
                all_chunks.append((meta, chunk[0]))
        if "embedding" in meta and meta["embedding"]:
            embedding_cache[file_path] = np.array(meta["embedding"])
    progress_bar.progress((idx + 1) / len(all_files), text=f"Indexed {idx+1}/{len(all_files)} files")
progress_bar.empty()
save_global_aliases(updated_global_aliases)

# Only show alias review UI if new columns were mapped
if new_columns_mapped:
    if st.checkbox("🔧 Edit Column Alias Mappings"):
        editable_aliases = st.text_area(
            "Edit mapping as JSON if needed:",
            value=json.dumps(updated_global_aliases, indent=2),
            key="alias_editor"
        )
        if st.button("💾 Save Updated Aliases"):
            try:
                aliases = json.loads(editable_aliases)
                save_global_aliases(aliases)
                st.success("Aliases updated successfully.")
            except Exception:
                st.warning("Failed to parse and save aliases.")

with st.form("question_form", clear_on_submit=False):

    user_query = st.text_input("Your question:", value=st.session_state.get("last_query", ""))
    submit = st.form_submit_button("Ask")

if submit and user_query.strip():
    st.info("Processing your question. This may take a moment...")
    scored_chunks = []
    query_embedding = get_embedding(user_query)
    for meta, chunk in all_chunks:
        meta_text = " ".join([
            str(meta.get("title", "")).lower(),
            str(meta.get("category", "")).lower(),
            " ".join(meta.get("tags", [])),
            str(meta.get("source_file", "")).lower(),
            str(meta.get("summary", "")).lower(),
            chunk.lower()
        ])
        keyword_score = sum(word in meta_text for word in user_query.lower().split())
        emb = embedding_cache.get(meta.get("source_file"))
        semantic_score = cosine_similarity(query_embedding, emb) if emb is not None else 0.0
        hybrid_score = 0.5 * keyword_score + 0.5 * semantic_score
        scored_chunks.append((hybrid_score, meta, chunk, keyword_score, semantic_score))

    scored_chunks.sort(reverse=True, key=lambda x: x[0])
    top_chunks = scored_chunks[:3] if scored_chunks else []

    context = "\n\n".join([chunk for _, _, chunk, _, _ in top_chunks])
    sources = [meta.get("source_file", "") for _, meta, _, _, _ in top_chunks]
    scores = [{"file": meta.get("source_file", ""), "score": score, "keyword": kw, "semantic": sem}
              for score, meta, _, kw, sem in top_chunks]

    # --- Enhanced Excel Q&A for initial question ---
    top_meta = top_chunks[0][1] if top_chunks else None
    top_file = top_meta.get("source_file") if top_meta else None
    if top_file and top_file.lower().endswith('.xlsx'):
        column_aliases = top_meta.get("column_aliases", {})
        excel_qa(top_file, user_query, column_aliases)
    else:
        # Fallback: original LLM reasoning answer for non-Excel files
        system_prompt = (
            "You are a business analyst answering questions using internal documents (Excel, Word, PDF, PPTX). "
            "For each question, follow this reasoning chain:\n"
            "1. Identify the key data needed to answer the question.\n"
            "2. Retrieve or summarize the relevant information from the provided context.\n"
            "3. If the answer is quantitative and a chart would help, describe the chart (or provide code for matplotlib/seaborn if possible).\n"
            "4. Explain what the result or chart shows.\n"
            "5. Suggest possible causes or business insights that could explain the observed pattern, using supporting evidence from other loaded data if possible.\n"
            "6. Cite the source file(s) used for the answer.\n"
            "7. Rate your confidence in the answer (0-100) based on relevance scores and verification.\n"
            "Use only the provided context and cite sources."
        )
        try:
            response = client.chat.completions.create(
                model=OPENAI_CHAT_MODEL,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": f"Context:\n{context}\n\nQuestion: {user_query}"}
                ],
                temperature=0.3
            )
            answer = response.choices[0].message.content.strip()
        except Exception as e:
            answer = f"LLM answer failed: {e}"
        st.markdown(f"**Answer:**\n\n{answer}")
        st.markdown("**Sources used:**")
        for src in sources:
            st.write(f"- {src}")
        st.markdown("**Relevance scores:**")
        for s in scores:
            st.write(f"{s['file']}: hybrid={s['score']:.2f}, keyword={s['keyword']}, semantic={s['semantic']:.2f}")

    verification = verify_answer(
        st.session_state.get("last_answer_table", "") if top_file and top_file.lower().endswith('.xlsx') else answer,
        context,
        user_query
    )
    root_cause_analysis = mine_for_root_causes(user_query, all_chunks, top_chunks)

    st.markdown("### \U0001F4CA Predictive Modeling Options")
    predictive_model_result = None

    if st.button("\u2699\ufe0f Run Prebuilt Model"):
        predictive_model_result = predictive_modeling_prebuilt(user_query, top_chunks)
        st.write(predictive_model_result)

    if st.button("\U0001F6E0 Build and Train Model (LLM-guided)"):
        with st.spinner("Generating and training your model..."):
            result = build_and_run_model(user_query, top_chunks)
            st.write(result)

    st.markdown("**Verification:**")
    st.write(verification)
    st.markdown("**Root Cause/Data Mining Analysis:**")
    st.write(root_cause_analysis)

    st.session_state.last_query = user_query
    st.session_state.chat_history.append({"role": "user", "content": user_query})
    st.session_state.chat_history.append({"role": "assistant", "content": st.session_state.get("last_answer_table", answer)})
    st.session_state.query_log.append({
        "timestamp": datetime.now().isoformat(),
        "question": user_query,
        "files_used": sources,
        "scores": scores,
        "answer": st.session_state.get("last_answer_table", answer),
        "verification": verification,
        "root_cause_analysis": root_cause_analysis,
        "predictive_model_result": predictive_model_result if predictive_model_result else "",
    })

else:
    st.write("Ask a question to get started.")

with st.expander("Show Query Log"):
    for entry in st.session_state.query_log:
        st.write(entry)

with st.expander("Show Chat History"):
    for msg in st.session_state.chat_history:
        st.write(f"{msg['role'].capitalize()}: {msg['content']}")

# --- Download alias file ---
if os.path.exists(GLOBAL_ALIAS_PATH):
    with open(GLOBAL_ALIAS_PATH, "r", encoding="utf-8") as f:
        global_alias_json = f.read()
else:
    global_alias_json = "{}"
st.download_button(
    label="Download global_column_aliases.json",
    data=global_alias_json,
    file_name="global_column_aliases.json",
    mime="application/json"
)
