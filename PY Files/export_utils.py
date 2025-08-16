# PY Files/export_utils.py
"""
Export utilities for SCIE Ethos Streamlit UI.
Handles XLSX, PDF, MD, PPT, DOCX exports with Dropbox + S3 integration.
"""

import os
import json
import time
import io
from datetime import datetime
from typing import Dict, Any, List, Optional, Union
from pathlib import Path

import pandas as pd
import streamlit as st

# Try to import optional dependencies
try:
    from docx import Document
    from docx.shared import Inches
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPT_AVAILABLE = True
except ImportError:
    PPT_AVAILABLE = False

try:
    import boto3
    S3_AVAILABLE = True
except ImportError:
    S3_AVAILABLE = False

# Local imports
try:
    from dbx_utils import upload_bytes, upload_json
    from path_utils import get_project_paths
    DBX_AVAILABLE = True
except ImportError:
    DBX_AVAILABLE = False

class ExportManager:
    """Manages exports to Dropbox and S3 with multiple formats."""
    
    def __init__(self, service_level: float = 0.95):
        self.service_level = service_level
        self.export_timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        # Get project paths
        try:
            self.paths = get_project_paths()
        except:
            self.paths = None
    
    def export_to_xlsx(self, data: Dict[str, Any], filename: str = None) -> bytes:
        """Export data to XLSX format."""
        if filename is None:
            filename = f"export_{self.export_timestamp}.xlsx"
        
        # Create Excel writer
        output = io.BytesIO()
        with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
            # Export conversation messages
            if 'messages' in data:
                messages_df = pd.DataFrame(data['messages'])
                messages_df.to_excel(writer, sheet_name='Conversation', index=False)
            
            # Export metadata
            if 'metadata' in data:
                metadata_df = pd.DataFrame([data['metadata']])
                metadata_df.to_excel(writer, sheet_name='Metadata', index=False)
            
            # Export sources
            if 'sources' in data:
                sources_df = pd.DataFrame(data['sources'])
                sources_df.to_excel(writer, sheet_name='Sources', index=False)
            
            # Export confidence history
            if 'confidence_history' in data:
                conf_df = pd.DataFrame(data['confidence_history'], columns=['confidence_score'])
                conf_df.to_excel(writer, sheet_name='Confidence', index=False)
        
        output.seek(0)
        return output.read()
    
    def export_to_markdown(self, data: Dict[str, Any], filename: str = None) -> str:
        """Export data to Markdown format."""
        if filename is None:
            filename = f"export_{self.export_timestamp}.md"
        
        lines = []
        lines.append(f"# SCIE Ethos Export - {filename}")
        lines.append(f"**Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        lines.append(f"**Service Level:** {self.service_level:.1%}")
        lines.append("")
        
        # Export conversation
        if 'messages' in data:
            lines.append("## Conversation")
            lines.append("")
            for msg in data['messages']:
                role = msg.get('role', 'unknown').title()
                content = msg.get('content', '')
                lines.append(f"### {role}")
                lines.append("")
                lines.append(content)
                lines.append("")
        
        # Export metadata
        if 'metadata' in data:
            lines.append("## Metadata")
            lines.append("")
            for key, value in data['metadata'].items():
                lines.append(f"**{key}:** {value}")
            lines.append("")
        
        # Export sources
        if 'sources' in data:
            lines.append("## Sources")
            lines.append("")
            for source in data['sources']:
                lines.append(f"- {source}")
            lines.append("")
        
        return "\n".join(lines)
    
    def export_to_docx(self, data: Dict[str, Any], filename: str = None) -> bytes:
        """Export data to DOCX format."""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx not available")
        
        if filename is None:
            filename = f"export_{self.export_timestamp}.docx"
        
        doc = Document()
        
        # Title
        title = doc.add_heading(f'SCIE Ethos Export - {filename}', 0)
        
        # Metadata
        doc.add_heading('Metadata', level=1)
        doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        doc.add_paragraph(f"Service Level: {self.service_level:.1%}")
        
        # Conversation
        if 'messages' in data:
            doc.add_heading('Conversation', level=1)
            for msg in data['messages']:
                role = msg.get('role', 'unknown').title()
                content = msg.get('content', '')
                doc.add_heading(role, level=2)
                doc.add_paragraph(content)
        
        # Sources
        if 'sources' in data:
            doc.add_heading('Sources', level=1)
            for source in data['sources']:
                doc.add_paragraph(source, style='List Bullet')
        
        # Save to bytes
        output = io.BytesIO()
        doc.save(output)
        output.seek(0)
        return output.read()
    
    def export_to_pptx(self, data: Dict[str, Any], filename: str = None) -> bytes:
        """Export data to PPTX format."""
        if not PPT_AVAILABLE:
            raise ImportError("python-pptx not available")
        
        if filename is None:
            filename = f"export_{self.export_timestamp}.pptx"
        
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = f"SCIE Ethos Export"
        subtitle.text = f"{filename}\n{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Metadata slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Metadata"
        content = slide.placeholders[1]
        content.text = f"Service Level: {self.service_level:.1%}\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Conversation slides
        if 'messages' in data:
            for i, msg in enumerate(data['messages']):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                role = msg.get('role', 'unknown').title()
                content_text = msg.get('content', '')[:500] + "..." if len(msg.get('content', '')) > 500 else msg.get('content', '')
                slide.shapes.title.text = f"{role} - Message {i+1}"
                content = slide.placeholders[1]
                content.text = content_text
        
        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()
    
    def save_to_dropbox(self, content: Union[bytes, str], filename: str, file_type: str = "xlsx") -> bool:
        """Save export to Dropbox."""
        if not DBX_AVAILABLE:
            st.warning("Dropbox integration not available")
            return False
        
        try:
            # Determine export path
            if self.paths and hasattr(self.paths, 'exports_folder'):
                export_path = f"{self.paths.exports_folder}/{filename}"
            else:
                export_path = f"Project_Root/05_Exports/{filename}"
            
            # Upload based on content type
            if isinstance(content, str):
                # For markdown, convert to bytes
                content_bytes = content.encode('utf-8')
                upload_bytes(content_bytes, export_path)
            else:
                # For binary content
                upload_bytes(content, export_path)
            
            st.success(f"✅ Exported to Dropbox: {export_path}")
            return True
            
        except Exception as e:
            st.error(f"❌ Dropbox export failed: {e}")
            return False
    
    def save_to_s3(self, content: Union[bytes, str], filename: str, file_type: str = "xlsx") -> bool:
        """Save export to S3 for audit purposes."""
        if not S3_AVAILABLE:
            st.warning("S3 integration not available")
            return False
        
        try:
            # Get S3 credentials from Streamlit secrets
            s3_client = boto3.client(
                's3',
                region_name=st.secrets.get("AWS_DEFAULT_REGION"),
                aws_access_key_id=st.secrets.get("AWS_ACCESS_KEY_ID"),
                aws_secret_access_key=st.secrets.get("AWS_SECRET_ACCESS_KEY"),
            )
            
            bucket = st.secrets.get("S3_BUCKET")
            prefix = st.secrets.get("S3_PREFIX", "").rstrip("/")
            
            # Create S3 key
            s3_key = f"{prefix}/exports/{self.export_timestamp}/{filename}"
            
            # Determine content type
            content_types = {
                "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "md": "text/markdown",
                "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            }
            
            content_type = content_types.get(file_type, "application/octet-stream")
            
            # Upload to S3
            if isinstance(content, str):
                content_bytes = content.encode('utf-8')
            else:
                content_bytes = content
            
            s3_client.put_object(
                Bucket=bucket,
                Key=s3_key,
                Body=content_bytes,
                ContentType=content_type
            )
            
            st.success(f"✅ Exported to S3: s3://{bucket}/{s3_key}")
            return True
            
        except Exception as e:
            st.error(f"❌ S3 export failed: {e}")
            return False
    
    def export_all_formats(self, data: Dict[str, Any], base_filename: str = None) -> Dict[str, bool]:
        """Export data to all available formats."""
        if base_filename is None:
            base_filename = f"scie_ethos_export_{self.export_timestamp}"
        
        results = {}
        
        # XLSX export
        try:
            xlsx_content = self.export_to_xlsx(data, f"{base_filename}.xlsx")
            results['xlsx'] = self.save_to_dropbox(xlsx_content, f"{base_filename}.xlsx", "xlsx")
            if results['xlsx']:
                self.save_to_s3(xlsx_content, f"{base_filename}.xlsx", "xlsx")
        except Exception as e:
            st.error(f"XLSX export failed: {e}")
            results['xlsx'] = False
        
        # Markdown export
        try:
            md_content = self.export_to_markdown(data, f"{base_filename}.md")
            results['md'] = self.save_to_dropbox(md_content, f"{base_filename}.md", "md")
            if results['md']:
                self.save_to_s3(md_content, f"{base_filename}.md", "md")
        except Exception as e:
            st.error(f"Markdown export failed: {e}")
            results['md'] = False
        
        # DOCX export
        if DOCX_AVAILABLE:
            try:
                docx_content = self.export_to_docx(data, f"{base_filename}.docx")
                results['docx'] = self.save_to_dropbox(docx_content, f"{base_filename}.docx", "docx")
                if results['docx']:
                    self.save_to_s3(docx_content, f"{base_filename}.docx", "docx")
            except Exception as e:
                st.error(f"DOCX export failed: {e}")
                results['docx'] = False
        else:
            results['docx'] = False
        
        # PPTX export
        if PPT_AVAILABLE:
            try:
                pptx_content = self.export_to_pptx(data, f"{base_filename}.pptx")
                results['pptx'] = self.save_to_dropbox(pptx_content, f"{base_filename}.pptx", "pptx")
                if results['pptx']:
                    self.save_to_s3(pptx_content, f"{base_filename}.pptx", "pptx")
            except Exception as e:
                st.error(f"PPTX export failed: {e}")
                results['pptx'] = False
        else:
            results['pptx'] = False
        
        return results
