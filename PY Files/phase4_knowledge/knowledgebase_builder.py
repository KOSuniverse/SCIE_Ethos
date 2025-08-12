# PY Files/phase4_knowledge/knowledgebase_builder.py
# Phase 4 — Knowledge Base Builder (final)
# Ingests supported docs, chunks, embeds with OpenAI, stores in FAISS.
# Key loading order: streamlit.secrets["OPENAI_API_KEY"] -> os.environ["OPENAI_API_KEY"]

from __future__ import annotations
import os, re, gc, json, time, argparse, pickle, hashlib
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import List, Dict, Optional, Tuple

# ---------- Optional extractors (guarded) ----------
try: import pdfplumber
except Exception: pdfplumber = None
try:
    from PIL import Image
except Exception:
    Image = None
try: import pytesseract
except Exception: pytesseract = None
try: import docx  # python-docx
except Exception: docx = None
try:
    from pptx import Presentation  # python-pptx
except Exception:
    Presentation = None
try: import openpyxl
except Exception: openpyxl = None

# ---------- Required libs ----------
import numpy as np
import faiss
try: import tiktoken
except Exception: tiktoken = None

# OpenAI client
from openai import OpenAI

# Enterprise foundation imports
try:
    from constants import PROJECT_ROOT
    from path_utils import join_root, canon_path
except ImportError:
    # Fallback for standalone usage
    PROJECT_ROOT = "/Project_Root"
    def join_root(path): return f"{PROJECT_ROOT}/{path.strip('/')}"
    def canon_path(path): return str(Path(path).resolve())

# Cloud storage operations (optional)
try:
    from .dbx_utils import upload_bytes, read_file_bytes, upload_json
    DBX_AVAILABLE = True
except ImportError:
    try:
        # Try importing from parent directory
        import sys
        import os
        sys.path.append(os.path.dirname(os.path.dirname(__file__)))
        from dbx_utils import upload_bytes, read_file_bytes, upload_json
        DBX_AVAILABLE = True
    except ImportError:
        DBX_AVAILABLE = False

# ---------- Relative KB paths (enterprise standard) ----------
KB_SUBDIR     = "06_LLM_Knowledge_Base"
INDEX_REL     = f"{KB_SUBDIR}/document_index.faiss"
DOCSTORE_REL  = f"{KB_SUBDIR}/docstore.pkl"
MANIFEST_REL  = f"{KB_SUBDIR}/manifest.json"
CHUNKS_REL    = f"{KB_SUBDIR}/chunks"

# ---------- Config ----------
EMBED_MODEL = os.environ.get("EMBED_MODEL", "text-embedding-3-small")
CHUNK_TOKENS = int(os.environ.get("KB_CHUNK_TOKENS", "700"))
CHUNK_OVERLAP_TOKENS = int(os.environ.get("KB_CHUNK_OVERLAP_TOKENS", "100"))
TOKENS_PER_CHAR_APPROX = 0.25
MIN_TEXT_CHARS_PER_PAGE = 40

SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".txt", ".md"}
TEXT_EXTS = {".txt", ".md"}
# --- Supported source extensions ---
ALLOWED_EXTS = {".pdf", ".docx", ".pptx", ".txt", ".md"}
DEFAULT_SCAN_FOLDERS = ["06_LLM_Knowledge_Base"]

CHUNKS_REL = f"{KB_SUBDIR}/chunks"

# ---------- Data ----------
@dataclass
class FileRecord:
    path: str
    mtime: float
    size: int
    sha256: str
    ext: str
    pages: Optional[int] = None
    embedded_chunks: int = 0

@dataclass
class Chunk:
    doc_id: str
    file_path: str
    source_type: str
    page_range: Optional[Tuple[int, int]]
    text: str
    tokens_est: int

# ---------- Utils ----------
def ensure_dirs(project_root: Path):
    """Create directory structure. Handles both local and cloud paths."""
    # For cloud paths (starting with /), we need to use Dropbox API calls
    # For local paths, use regular filesystem operations
    project_root_str = str(project_root)
    
    if project_root_str.startswith("/") and DBX_AVAILABLE:
        # This is a Dropbox path - directories are created implicitly when files are uploaded
        # We'll create a placeholder file to ensure the directories exist
        try:
            # Create placeholder files to ensure directories exist
            kb_path = f"{project_root_str.rstrip('/')}/{KB_SUBDIR}"
            chunks_path = f"{kb_path}/chunks"
            
            # Create .gitkeep files to ensure directories exist
            upload_bytes(f"{kb_path}/.gitkeep", b"", mode="overwrite")
            upload_bytes(f"{chunks_path}/.gitkeep", b"", mode="overwrite")
        except Exception:
            # If cloud operations fail, fall back to local operations
            # This handles cases where dbx_utils isn't available or configured
            (project_root / "06_LLM_Knowledge_Base").mkdir(parents=True, exist_ok=True)
            (project_root / "06_LLM_Knowledge_Base" / "chunks").mkdir(parents=True, exist_ok=True)
    else:
        # Local filesystem operations
        (project_root / "06_LLM_Knowledge_Base").mkdir(parents=True, exist_ok=True)
        (project_root / "06_LLM_Knowledge_Base" / "chunks").mkdir(parents=True, exist_ok=True)

def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for b in iter(lambda: f.read(1024 * 1024), b""): h.update(b)
    return h.hexdigest()

def list_source_files(project_root: Path, scan_folders: List[str]) -> List[Path]:
    """List source files from local or cloud storage."""
    project_root_str = str(project_root)
    files: List[Path] = []
    
    # For cloud paths, we can't scan directories in the same way
    # Return empty list for now - files will need to be added explicitly
    if project_root_str.startswith("/"):
        # Cloud mode - source files would need to be listed differently
        # For now, return empty list since we can't scan cloud directories
        # Files can be added explicitly via the CLI add command
        return []
    
    # Local filesystem scanning
    for rel in scan_folders:
        folder = (project_root / rel).resolve()
        if not folder.exists(): continue
        for ext in SUPPORTED_EXTS:
            files.extend(folder.rglob(f"*{ext}"))
    return sorted(set(files))

def load_manifest(path: Path) -> Dict[str, FileRecord]:
    """Load manifest from local or cloud storage."""
    path_str = str(path)
    
    if path_str.startswith("/") and DBX_AVAILABLE:
        # Cloud path - use Dropbox API
        try:
            data = json.loads(read_file_bytes(path_str).decode("utf-8"))
            return {k: FileRecord(**v) for k, v in data.items()}
        except Exception:
            # File doesn't exist or other error
            return {}
    else:
        # Local filesystem
        if path.exists():
            data = json.loads(path.read_text(encoding="utf-8"))
            return {k: FileRecord(**v) for k, v in data.items()}
        return {}

def save_manifest(path: Path, manifest: Dict[str, FileRecord]):
    """Save manifest to local or cloud storage."""
    path_str = str(path)
    data = json.dumps({k: asdict(v) for k, v in manifest.items()}, indent=2)
    
    if path_str.startswith("/") and DBX_AVAILABLE:
        # Cloud path - use Dropbox API
        try:
            upload_bytes(path_str, data.encode("utf-8"), mode="overwrite")
        except Exception:
            # Fall back to local operations
            path.write_text(data, encoding="utf-8")
    else:
        # Local filesystem
        path.write_text(data, encoding="utf-8")

def estimate_tokens(text: str) -> int:
    if not text: return 0
    if tiktoken:
        try: enc = tiktoken.get_encoding("cl100k_base")
        except Exception: enc = tiktoken.encoding_for_model("gpt-4o")
        return len(enc.encode(text))
    return int(len(text) * TOKENS_PER_CHAR_APPROX)

def chunk_text(text: str, chunk_tokens=CHUNK_TOKENS, overlap_tokens=CHUNK_OVERLAP_TOKENS) -> List[str]:
    if not text: return []
    if tiktoken:
        try: enc = tiktoken.get_encoding("cl100k_base")
        except Exception: enc = tiktoken.encoding_for_model("gpt-4o")
        toks = enc.encode(text)
        out, start, step = [], 0, max(1, chunk_tokens - overlap_tokens)
        while start < len(toks):
            end = min(len(toks), start + chunk_tokens)
            out.append(enc.decode(toks[start:end]))
            start += step
        return out
    approx_chars = int(chunk_tokens / TOKENS_PER_CHAR_APPROX)
    approx_overlap = int(overlap_tokens / TOKENS_PER_CHAR_APPROX)
    out, start, step = [], 0, max(1, approx_chars - approx_overlap)
    while start < len(text):
        out.append(text[start:min(len(text), start + approx_chars)])
        start += step
    return out

def safe_text(s: Optional[str]) -> str:
    if not s: return ""
    s = re.sub(r"\s+\n", "\n", s)
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()

def detect_source_type(path: Path) -> str:
    return {
        ".pdf": "pdf", ".docx": "docx", ".pptx": "pptx",
        ".xlsx": "xlsx", ".xls": "xlsx", ".txt": "txt", ".md": "txt"
    }.get(path.suffix.lower(), "txt")

# ---------- Extractors ----------
def extract_pdf_text(path: Path) -> Tuple[List[str], int]:
    if pdfplumber is None:
        raise RuntimeError("Missing pdfplumber (pip install pdfplumber)")
    pages = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            txt = safe_text(page.extract_text() or "")
            # selective OCR if page is mostly empty and OCR deps exist
            if len(txt) < MIN_TEXT_CHARS_PER_PAGE and pytesseract and Image:
                try:
                    img = page.to_image(resolution=300).original
                    ocr = pytesseract.image_to_string(img)
                    if ocr: txt = safe_text(txt + "\n\n" + ocr)
                except Exception:
                    pass
            pages.append(txt)
    return pages, len(pages)

def extract_docx_text(path: Path) -> str:
    if docx is None: raise RuntimeError("Missing python-docx")
    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs]
    for t in d.tables:
        for r in t.rows:
            parts.append("\t".join([c.text for c in r.cells]))
    return safe_text("\n".join(parts))

def extract_pptx_text(path: Path) -> str:
    if Presentation is None: raise RuntimeError("Missing python-pptx")
    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"): parts.append(shape.text)
    return safe_text("\n".join(parts))

def extract_xlsx_text(path: Path) -> str:
    if openpyxl is None: raise RuntimeError("Missing openpyxl")
    wb = openpyxl.load_workbook(filename=str(path), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        row_count = 0
        for row in ws.iter_rows(values_only=True):
            parts.append(" | ".join([str(c) if c is not None else "" for c in row[:50]]))
            row_count += 1
            if row_count >= 5000:
                parts.append("... [truncated rows]")
                break
    wb.close()
    return safe_text("\n".join(parts))

def extract_txt(path: Path) -> str:
    for enc in ("utf-8", "latin-1"):
        try: return safe_text(path.read_text(encoding=enc, errors="ignore"))
        except Exception: continue
    return ""

def extract_file(path: Path) -> Tuple[str, Optional[List[str]], Optional[int], str]:
    stype = detect_source_type(path)
    if stype == "pdf":
        per_page, n = extract_pdf_text(path)
        return "", per_page, n, stype
    if stype == "docx": return extract_docx_text(path), None, None, stype
    if stype == "pptx": return extract_pptx_text(path), None, None, stype
    if stype == "xlsx": return extract_xlsx_text(path), None, None, stype
    return extract_txt(path), None, None, "txt"

# ---------- Chunk builders ----------
def chunks_from_pdf(per_page: List[str], file_path: Path, doc_id: str) -> List[Chunk]:
    chunks, buf, start_page, token_count = [], [], 1, 0
    for idx, page_text in enumerate(per_page, start=1):
        if not page_text: continue
        ptoks = estimate_tokens(page_text)
        if token_count + ptoks > CHUNK_TOKENS and buf:
            text = safe_text("\n\n".join(buf))
            for part in chunk_text(text):
                chunks.append(Chunk(doc_id, str(file_path), "pdf", (start_page, idx - 1), part, estimate_tokens(part)))
            buf, token_count, start_page = [], 0, idx
        buf.append(page_text); token_count += ptoks
    if buf:
        text = safe_text("\n\n".join(buf))
        for part in chunk_text(text):
            chunks.append(Chunk(doc_id, str(file_path), "pdf", (start_page, len(per_page)), part, estimate_tokens(part)))
    return chunks

def chunks_from_flat(text: str, file_path: Path, doc_id: str, source_type: str) -> List[Chunk]:
    return [Chunk(doc_id, str(file_path), source_type, None, part, estimate_tokens(part)) for part in chunk_text(text)]

# ---------- Key loader ----------
def _load_openai_key() -> str:
    # Try Streamlit secrets first
    try:
        import streamlit as st  # only if running inside Streamlit
        key = st.secrets.get("OPENAI_API_KEY")
        if key: return key.strip()
    except Exception:
        pass
    # Fallback to env var
    key = os.environ.get("OPENAI_API_KEY")
    if key: return key.strip()
    try:
        from dotenv import load_dotenv
        load_dotenv()
        key = os.environ.get("OPENAI_API_KEY")
        if key:
            return key.strip()
    except Exception:
        pass
    raise RuntimeError(
        "OPENAI_API_KEY not found. Set it in Streamlit secrets, environment, or a .env file."
    )

# ---------- Embeddings / FAISS ----------
def get_openai_client() -> OpenAI:
    key = _load_openai_key()
    return OpenAI(api_key=key)

def embed_texts(client: OpenAI, texts: List[str], model: str = EMBED_MODEL, batch_size: int = 96) -> np.ndarray:
    vectors: List[List[float]] = []
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i + batch_size]
        resp = client.embeddings.create(model=model, input=batch)
        for d in resp.data: vectors.append(d.embedding)
        time.sleep(0.02)  # tiny pacing
    return np.array(vectors, dtype="float32")

def write_faiss_index(path: Path, index: faiss.IndexFlatIP): 
    """Write FAISS index to local or cloud storage."""
    path_str = str(path)
    
    if path_str.startswith("/") and DBX_AVAILABLE:
        # Cloud path - write to temp file then upload
        try:
            import tempfile
            
            with tempfile.NamedTemporaryFile(delete=False) as tmp:
                faiss.write_index(index, tmp.name)
                tmp.flush()
                
                with open(tmp.name, "rb") as f:
                    data = f.read()
                upload_bytes(path_str, data, mode="overwrite")
                
                os.unlink(tmp.name)
        except Exception:
            # Fall back to local operations
            faiss.write_index(index, str(path))
    else:
        # Local filesystem
        faiss.write_index(index, str(path))

def read_faiss_index(path: Path) -> faiss.IndexFlatIP: 
    """Read FAISS index from local or cloud storage."""
    path_str = str(path)
    
    if path_str.startswith("/") and DBX_AVAILABLE:
        # Cloud path - download to temp file then read
        try:
            import tempfile
            
            data = read_file_bytes(path_str)
            with tempfile.NamedTemporaryFile(delete=False) as tmp:
                tmp.write(data)
                tmp.flush()
                
                index = faiss.read_index(tmp.name)
                os.unlink(tmp.name)
                return index
        except Exception:
            # Fall back to local operations or raise
            if Path(path_str).exists():
                return faiss.read_index(str(path))
            raise
    else:
        # Local filesystem
        return faiss.read_index(str(path))

def load_docstore(path: Path) -> Dict[str, list]:
    """Load docstore from local or cloud storage."""
    path_str = str(path)
    
    if path_str.startswith("/") and DBX_AVAILABLE:
        # Cloud path - use Dropbox API
        try:
            data = read_file_bytes(path_str)
            return pickle.loads(data)
        except Exception:
            # File doesn't exist or other error
            return {}
    else:
        # Local filesystem
        if path.exists():
            with path.open("rb") as f: 
                return pickle.load(f)
        return {}

def save_docstore(path: Path, docstore: Dict[str, list]):
    """Save docstore to local or cloud storage."""
    path_str = str(path)
    
    if path_str.startswith("/") and DBX_AVAILABLE:
        # Cloud path - use Dropbox API
        try:
            data = pickle.dumps(docstore)
            upload_bytes(path_str, data, mode="overwrite")
        except Exception:
            # Fall back to local operations
            with path.open("wb") as f: 
                pickle.dump(docstore, f)
    else:
        # Local filesystem
        with path.open("wb") as f: 
            pickle.dump(docstore, f)

def save_chunk_records(chunks_dir: Path, doc_id: str, chunks: List[Chunk]):
    """Save chunk records to local or cloud storage."""
    chunks_dir_str = str(chunks_dir)
    file_path = f"{chunks_dir_str.rstrip('/')}/{doc_id}.jsonl"
    
    # Build JSONL content
    lines = []
    for c in chunks:
        lines.append(json.dumps({
            "doc_id": c.doc_id, "file_path": c.file_path, "source_type": c.source_type,
            "page_range": c.page_range, "tokens_est": c.tokens_est, "text": c.text
        }))
    content = "\n".join(lines) + "\n"
    
    if chunks_dir_str.startswith("/") and DBX_AVAILABLE:
        # Cloud path - use Dropbox API
        try:
            upload_bytes(file_path, content.encode("utf-8"), mode="overwrite")
        except Exception:
            # Fall back to local operations
            out = chunks_dir / f"{doc_id}.jsonl"
            with out.open("w", encoding="utf-8") as f:
                f.write(content)
    else:
        # Local filesystem
        out = chunks_dir / f"{doc_id}.jsonl"
        with out.open("w", encoding="utf-8") as f:
            f.write(content)

def normalize_vectors_for_ip(vectors: np.ndarray) -> np.ndarray:
    norms = np.linalg.norm(vectors, axis=1, keepdims=True); norms[norms == 0] = 1.0
    return vectors / norms

def cloud_file_exists(path_str: str) -> bool:
    """Check if a file exists in cloud storage."""
    if not DBX_AVAILABLE:
        return False
    try:
        read_file_bytes(path_str)
        return True
    except Exception:
        return False

def file_exists_any(path: Path) -> bool:
    """Check if file exists in local or cloud storage."""
    path_str = str(path)
    if path_str.startswith("/") and DBX_AVAILABLE:
        return cloud_file_exists(path_str)
    else:
        return path.exists()

# ---------- Build / Update ----------
def scan_and_hash(project_root: Path, scan_folders: List[str]) -> Dict[str, FileRecord]:
    files = list_source_files(project_root, scan_folders)
    out: Dict[str, FileRecord] = {}
    for p in files:
        try:
            stat = p.stat()
            out[str(p)] = FileRecord(str(p), stat.st_mtime, stat.st_size, sha256_file(p), p.suffix.lower())
        except Exception:
            continue
    return out

def need_embedding(manifest: Dict[str, FileRecord], rec: FileRecord) -> bool:
    ex = manifest.get(rec.path)
    return (not ex) or (ex.sha256 != rec.sha256)

def build_or_update_knowledgebase(
    project_root: str,
    scan_folders: Optional[List[str]] = None,
    force_rebuild: bool = False,
    include_text_files: bool = True,
) -> Dict[str, any]:
    # Don't resolve cloud paths to avoid converting to local filesystem
    if project_root.startswith("/"):
        root = Path(project_root)  # Keep as cloud path
    else:
        root = Path(project_root).resolve()  # Only resolve local paths
    
    ensure_dirs(root)

    scans = scan_folders if scan_folders else DEFAULT_SCAN_FOLDERS
    manifest_path, index_path = root / MANIFEST_REL, root / INDEX_REL
    docstore_path, chunks_dir = root / DOCSTORE_REL, root / CHUNKS_REL

    prev_manifest = load_manifest(manifest_path)
    docstore = load_docstore(docstore_path)
    current = scan_and_hash(root, scans)

    if not include_text_files:
        current = {k: v for k, v in current.items() if Path(k).suffix.lower() not in TEXT_EXTS}

    to_process = list(current.values()) if force_rebuild else [rec for rec in current.values() if need_embedding(prev_manifest, rec)]

    index: Optional[faiss.IndexFlatIP] = None
    if not force_rebuild and file_exists_any(index_path):
        try:
            index = read_faiss_index(index_path)
        except Exception:
            index = None
    
    client = get_openai_client()

    results: Dict[str, dict] = {}
    total_added = 0

    for rec in to_process:
        fp = Path(rec.path)
        rel = str(fp.relative_to(root)) if str(fp).startswith(str(root)) else rec.path

        try:
            flat_text, per_page, n_pages, stype = extract_file(fp)
        except Exception as e:
            results[rel] = {"status": "extract_error", "error": str(e)}
            continue

        if per_page is not None:
            chunks = chunks_from_pdf(per_page, fp, doc_id=rec.sha256[:16]); pages = n_pages
        else:
            chunks = chunks_from_flat(flat_text, fp, doc_id=rec.sha256[:16], source_type=stype); pages = None

        if not chunks:
            results[rel] = {"status": "no_chunks"}; continue

        save_chunk_records(chunks_dir, rec.sha256[:16], chunks)

        texts = [c.text for c in chunks]
        try:
            vecs = embed_texts(client, texts, model=EMBED_MODEL)
        except Exception as e:
            results[rel] = {"status": "embed_error", "error": str(e)}
            continue

        vecs = normalize_vectors_for_ip(vecs)
        if index is None: index = faiss.IndexFlatIP(vecs.shape[1])
        index.add(vecs); total_added += vecs.shape[0]

        # Store lightweight meta + texts aligned to vector order
        metas = []
        for c in chunks:
            metas.append({
                "doc_id": c.doc_id, "file_path": c.file_path, "source_type": c.source_type,
                "page_range": c.page_range, "tokens_est": c.tokens_est
            })
        docstore.setdefault("vectors", []).extend(metas)
        docstore.setdefault("texts", []).extend(texts)

        rec.pages = pages; rec.embedded_chunks = len(chunks)
        prev_manifest[rec.path] = rec
        results[rel] = {"status": "embedded", "chunks": len(chunks), "pages": pages}

        # periodic safe flush for large batches
        if total_added >= 2048:
            write_faiss_index(index_path, index)
            save_docstore(docstore_path, docstore)
            save_manifest(manifest_path, prev_manifest)
            total_added = 0
            gc.collect()

    # Final persist
    if index is not None: write_faiss_index(index_path, index)
    save_docstore(docstore_path, docstore)
    save_manifest(manifest_path, prev_manifest)

    return {
        "processed": len(to_process),
        "total_files_seen": len(current),
        "index_vectors_total": (0 if index is None else index.ntotal),
        "index_path": str(index_path),
        "docstore_path": str(docstore_path),
        "manifest_path": str(manifest_path),
        "results_by_file": results,
    }

# ---------- Status / helpers ----------
def status(project_root: str) -> Dict[str, any]:
    # Don't resolve cloud paths to avoid converting to local filesystem
    if project_root.startswith("/"):
        root = Path(project_root)  # Keep as cloud path
    else:
        root = Path(project_root).resolve()  # Only resolve local paths
    
    man = load_manifest(root / MANIFEST_REL)
    idx, ds, manp = root / INDEX_REL, root / DOCSTORE_REL, root / MANIFEST_REL
    return {
        "manifest_files": len(man),
        "index_exists": file_exists_any(idx),
        "docstore_exists": file_exists_any(ds),
        "manifest_exists": file_exists_any(manp),
        "index_path": str(idx),
        "docstore_path": str(ds),
        "manifest_path": str(manp),
        "project_root": project_root,
        "cloud_mode": str(root).startswith("/") and DBX_AVAILABLE,
        "dbx_available": DBX_AVAILABLE,
        "computed_project_root": PROJECT_ROOT,
    }

# ---------- CLI ----------
def _parse_args():
    ap = argparse.ArgumentParser(description="Phase 4 Knowledgebase Builder")
    ap.add_argument("--project-root", type=str, required=True, help="Path to Project_Root")
    ap.add_argument("--scan-folders", type=str, default=",".join(DEFAULT_SCAN_FOLDERS),
                    help="Comma-separated folders (relative to Project_Root)")
    ap.add_argument("--no-text", action="store_true", help="Exclude .txt/.md from ingestion")
    sub = ap.add_subparsers(dest="cmd", required=True)
    sub.add_parser("status", help="Show KB status")
    sub.add_parser("build", help="Incremental build/update")
    sub.add_parser("rebuild", help="Force full rebuild")
    p_add = sub.add_parser("add", help="Embed specific files (absolute or relative)")
    p_add.add_argument("--files", type=str, required=True, help="Comma-separated file paths")
    return ap.parse_args()

def _cli_add(root: Path, files_csv: str):
    manifest_path, index_path = root / MANIFEST_REL, root / INDEX_REL
    docstore_path, chunks_dir = root / DOCSTORE_REL, root / CHUNKS_REL
    ensure_dirs(root)
    prev_manifest = load_manifest(manifest_path)
    docstore = load_docstore(docstore_path)
    
    index = None
    if file_exists_any(index_path):
        try:
            index = read_faiss_index(index_path)
        except Exception:
            index = None
    
    client = get_openai_client()

    results, total_added = {}, 0
    for raw in [s.strip() for s in files_csv.split(",") if s.strip()]:
        p = Path(raw)
        # For cloud paths, don't resolve; for local paths, resolve relative to root
        if not str(root).startswith("/"):
            if not p.is_absolute(): 
                p = (root / p).resolve()
        
        # For cloud mode, we can't check .exists() - skip this check
        if not str(root).startswith("/") and not p.exists():
            results[raw] = {"status": "missing"}; continue
        try:
            stat = p.stat()
            rec = FileRecord(str(p), stat.st_mtime, stat.st_size, sha256_file(p), p.suffix.lower())
        except Exception as e:
            results[raw] = {"status": "hash_error", "error": str(e)}; continue

        try:
            flat_text, per_page, n_pages, stype = extract_file(p)
        except Exception as e:
            results[raw] = {"status": "extract_error", "error": str(e)}; continue

        if per_page is not None:
            chunks = chunks_from_pdf(per_page, p, doc_id=rec.sha256[:16]); pages = n_pages
        else:
            chunks = chunks_from_flat(flat_text, p, doc_id=rec.sha256[:16], source_type=stype); pages = None

        if not chunks:
            results[raw] = {"status": "no_chunks"}; continue

        save_chunk_records(chunks_dir, rec.sha256[:16], chunks)

        try:
            vecs = embed_texts(client, [c.text for c in chunks], model=EMBED_MODEL)
        except Exception as e:
            results[raw] = {"status": "embed_error", "error": str(e)}; continue

        vecs = normalize_vectors_for_ip(vecs)
        if index is None: index = faiss.IndexFlatIP(vecs.shape[1])
        index.add(vecs); total_added += vecs.shape[0]

        docstore.setdefault("vectors", []).extend([{
            "doc_id": c.doc_id, "file_path": c.file_path, "source_type": c.source_type,
            "page_range": c.page_range, "tokens_est": c.tokens_est
        } for c in chunks])
        docstore.setdefault("texts", []).extend([c.text for c in chunks])

        rec.pages = pages; rec.embedded_chunks = len(chunks)
        prev_manifest[rec.path] = rec
        results[raw] = {"status": "embedded", "chunks": len(chunks), "pages": pages}

        if total_added >= 2048:
            write_faiss_index(index_path, index); save_docstore(docstore_path, docstore); save_manifest(manifest_path, prev_manifest)
            total_added = 0; gc.collect()

    if index is not None: write_faiss_index(index_path, index)
    save_docstore(docstore_path, docstore); save_manifest(manifest_path, prev_manifest)
    print(json.dumps({"added": len(results), "results_by_file": results, "index_path": str(index_path)}, indent=2))

def main():
    args = _parse_args()
    # Don't resolve cloud paths to avoid converting to local filesystem
    if args.project_root.startswith("/"):
        root = Path(args.project_root)  # Keep as cloud path
    else:
        root = Path(args.project_root).resolve()  # Only resolve local paths
    
    scans = [s.strip() for s in args.scan_folders.split(",") if s.strip()]
    include_text = not args.no_text

    if args.cmd == "status":
        print(json.dumps(status(str(root)), indent=2)); return
    if args.cmd == "rebuild":
        print(json.dumps(build_or_update_knowledgebase(str(root), scans, True, include_text), indent=2)); return
    if args.cmd == "build":
        print(json.dumps(build_or_update_knowledgebase(str(root), scans, False, include_text), indent=2)); return
    if args.cmd == "add":
        _cli_add(root, args.files); return

if __name__ == "__main__":
    main()

