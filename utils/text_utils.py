# --- utils/text_utils.py ---

import nltk
import openpyxl
import pdfplumber
import pytesseract
import streamlit as st
from docx import Document
from pptx import Presentation
from PIL import Image

# This will be imported by the main file, so download_file should be passed as parameter
# or imported from the main scope

# Ensure NLTK data is available
try:
    nltk.data.find("tokenizers/punkt")
except LookupError:
    try:
        nltk.download("punkt")
    except:
        pass

try:
    nltk.data.find("tokenizers/punkt_tab")
except LookupError:
    try:
        nltk.download("punkt_tab")
    except:
        pass

def chunk_text(text, chunk_size=2000, overlap=200, max_chars=4000):
    """
    Splits text into sentence-based chunks with optional overlap and character limit.
    Returns a list of (chunk, start_idx, end_idx) tuples (matching original Jarvis format).
    """
    if not text or not text.strip():
        return []
    
    sentences = nltk.sent_tokenize(text)
    chunks = []
    current_chunk = []
    current_len = 0
    start_idx = 0
    
    for sent in sentences:
        sent = sent.strip()
        if not sent:
            continue
            
        if current_len + len(sent) > max_chars or len(current_chunk) >= chunk_size:
            if current_chunk:
                chunk_text = " ".join(current_chunk).strip()
                end_idx = start_idx + len(chunk_text)
                chunks.append((chunk_text, start_idx, end_idx))
                start_idx = end_idx
            
            # Overlap logic
            if overlap > 0 and len(current_chunk) > overlap:
                current_chunk = current_chunk[-overlap:]
                current_len = sum(len(s) for s in current_chunk)
            else:
                current_chunk = []
                current_len = 0
        
        current_chunk.append(sent)
        current_len += len(sent)
    
    # Add last chunk
    if current_chunk:
        chunk_text = " ".join(current_chunk).strip()
        end_idx = start_idx + len(chunk_text)
        chunks.append((chunk_text, start_idx, end_idx))
    
    return chunks

def extract_text_for_metadata(file_id, max_ocr_pages=5, download_file_func=None, get_drive_service_func=None):
    """
    Extracts human-readable content from supported document types.
    Returns:
        - For Excel (.xlsx): (text, sheet_names, columns_by_sheet)
        - For all others: (text, [], {})
    
    Args:
        file_id: Google Drive file ID
        max_ocr_pages: Maximum pages to process with OCR
        download_file_func: Function to download files from Google Drive
        get_drive_service_func: Function to get Google Drive service
    """
    if not download_file_func:
        raise ValueError("download_file_func must be provided")
    if not get_drive_service_func:
        raise ValueError("get_drive_service_func must be provided")
        
    try:
        # Download file from Google Drive using file_id
        file_stream = download_file_func(file_id)
        if not file_stream:
            return "", [], {}
        
        # Get file metadata to determine file type
        try:
            service = get_drive_service_func()
            file_metadata = service.files().get(fileId=file_id).execute()
            file_name = file_metadata.get('name', '')
            ext = file_name.lower().split('.')[-1] if '.' in file_name else ''
        except Exception:
            # Fallback: try to determine from content
            ext = ""
        
        # Try to determine file type and extract content accordingly
        if ext == "docx":
            file_stream.seek(0)
            doc = Document(file_stream)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return text, [], {}

        elif ext == "xlsx":
            file_stream.seek(0)
            wb = openpyxl.load_workbook(file_stream, read_only=True)
            return _extract_excel_content(wb)

        elif ext == "pptx":
            file_stream.seek(0)
            prs = Presentation(file_stream)
            text = "\n".join(
                shape.text for slide in prs.slides for shape in slide.shapes 
                if shape.has_text_frame
            )
            return text, [], {}

        elif ext == "pdf":
            file_stream.seek(0)
            return _extract_pdf_content(file_stream, max_ocr_pages)
        
        # Try to auto-detect if extension is unknown
        try:
            # Try Excel first
            file_stream.seek(0)
            wb = openpyxl.load_workbook(file_stream, read_only=True)
            return _extract_excel_content(wb)
        except:
            pass
        
        
        try:
            # Try Word document
            file_stream.seek(0)
            doc = Document(file_stream)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return text, [], {}
        except:
            pass
        
        try:
            # Try PowerPoint
            file_stream.seek(0)
            prs = Presentation(file_stream)
            text = "\n".join(
                shape.text for slide in prs.slides for shape in slide.shapes 
                if shape.has_text_frame
            )
            return text, [], {}
        except:
            pass
        
        try:
            # Try PDF
            file_stream.seek(0)
            return _extract_pdf_content(file_stream, max_ocr_pages)
        except:
            pass
        
    except Exception as e:
        st.warning(f"Failed to extract text from file: {e}")
    
    return "", [], {}

def _extract_excel_content(wb):
    """Extract content from Excel workbook."""
    text = []
    sheet_names = []
    columns_by_sheet = {}
    
    for sheet in wb.worksheets:
        sheet_names.append(sheet.title)
        text.append(f"Sheet: {sheet.title}")
        
        # Get headers
        headers = [cell.value for cell in next(sheet.iter_rows(max_row=1))]
        columns_by_sheet[sheet.title] = [str(h) for h in headers if h]
        
        if headers:
            text.append("Columns: " + " | ".join(columns_by_sheet[sheet.title]))
        
        # Get sample data (first few rows)
        for i, row in enumerate(sheet.iter_rows(min_row=2, max_row=6)):
            if i >= 5:  # Limit sample data
                break
            row_values = [str(cell.value) if cell.value is not None else "" for cell in row]
            if any(row_values):  # Only add non-empty rows
                text.append(" | ".join(row_values))
    
    return "\n".join(text), sheet_names, columns_by_sheet

def _extract_pdf_content(file_stream, max_ocr_pages=5):
    """Extract content from PDF with OCR fallback."""
    try:
        with pdfplumber.open(file_stream) as pdf:
            extracted_text = []
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    extracted_text.append(page_text)
            
            if any(extracted_text):
                return "\n".join(extracted_text), [], {}
            
            # No extractable text — try OCR fallback
            st.warning("No extractable text found. This PDF may require OCR.")
            st.info("OCR processing not available in cloud deployment (tesseract not installed).")
            return "", [], {}
            
    except Exception as e:
        st.warning(f"PDF extraction failed: {e}")
        return "", [], {}

def extract_structural_metadata(text, file_type):
    """Extract structural metadata from text based on file type."""
    if not text:
        return {}
    
    if file_type == "docx":
        return {
            "section_headings": [
                line.strip() for line in text.splitlines() 
                if line.strip().endswith(":") or line.strip().istitle()
            ]
        }
    elif file_type == "pdf":
        return {
            "section_headings": [
                line.strip() for line in text.splitlines() 
                if line.strip().endswith(":") or line.strip().istitle()
            ]
        }
    elif file_type == "pptx":
        return {
            "slide_titles": [
                line.strip() for line in text.splitlines() 
                if len(line.strip()) > 0
            ]
        }
    return {}
